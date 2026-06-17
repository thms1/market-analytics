"""
Telco Data Science Pipeline — Zambia Market Analytics
==========================================================
End-to-end data science roadmap with a step-by-step UI:

  Step 1  Data Ingestion          Upload CSV / use demo data
  Step 2  Exploratory Analysis    Stats, distributions, correlations
  Step 3  Data Cleaning           Nulls, duplicates, type coercion
  Step 4  Feature Engineering     Encoding, scaling, derived KPIs
  Step 5  Model Training          Six classifiers compared side-by-side
  Step 6  Model Evaluation        Accuracy, ROC-AUC, F1, confusion matrix
  Step 7  Bias Detection          Learning curves – overfit / underfit
  Step 8  Recommendations         Product & campaign segmentation
  Step 9  Data Story              Executive + technical PowerPoint exports
  Step 10 Market Intelligence     Zamtel vs MTN/Airtel KPI monitor & benchmarking

Run:
    py -3.12 -m streamlit run Data_Science_Pipeline.py
"""

from __future__ import annotations

import io
import json
import warnings
from datetime import datetime, timedelta
from pathlib import Path

import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import streamlit as st
from sklearn.ensemble import (
    AdaBoostClassifier,
    GradientBoostingClassifier,
    RandomForestClassifier,
)
from sklearn.impute import SimpleImputer
from sklearn.linear_model import LogisticRegression
from sklearn.metrics import (
    accuracy_score,
    classification_report,
    confusion_matrix,
    f1_score,
    precision_score,
    recall_score,
    roc_auc_score,
    roc_curve,
)
from sklearn.model_selection import learning_curve, train_test_split
from sklearn.neighbors import KNeighborsClassifier
from sklearn.preprocessing import LabelEncoder, StandardScaler
from sklearn.svm import SVC

warnings.filterwarnings("ignore")

# ─────────────────────────────────────────────────────────────────────────────
# PAGE CONFIG & STYLING
# ─────────────────────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="Zambia Telco Analytics",
    page_icon="📡",
    layout="wide",
    initial_sidebar_state="expanded",
)

st.markdown(
    """
<style>
    .step-header{font-size:1.6rem;font-weight:700;color:#1f4e79;
        border-bottom:3px solid #2196f3;padding-bottom:10px;margin-bottom:20px;}
    .kpi-card{background:linear-gradient(135deg,#1565c0,#42a5f5);
        border-radius:14px;padding:18px 12px;color:#fff;text-align:center;margin:4px 0;}
    .kpi-value{font-size:2.2rem;font-weight:800;line-height:1.1;}
    .kpi-label{font-size:.82rem;opacity:.88;margin-top:4px;}
    .kpi-card.green{background:linear-gradient(135deg,#2e7d32,#66bb6a);}
    .kpi-card.orange{background:linear-gradient(135deg,#e65100,#ffa726);}
    .kpi-card.purple{background:linear-gradient(135deg,#4a148c,#ab47bc);}
    .insight-box{background:#e3f2fd;border-left:5px solid #1976d2;
        padding:12px 16px;border-radius:4px;margin:8px 0;font-size:.95rem;}
    .warn-box{background:#fff8e1;border-left:5px solid #f57c00;
        padding:12px 16px;border-radius:4px;margin:8px 0;font-size:.95rem;}
    .ok-box{background:#e8f5e9;border-left:5px solid #388e3c;
        padding:12px 16px;border-radius:4px;margin:8px 0;font-size:.95rem;}
    .section-divider{border:none;border-top:2px dashed #90caf9;margin:20px 0;}
    .reco-card{border-radius:10px;padding:14px;margin:6px 0;color:#fff;font-size:.9rem;}
    .reco-red{background:linear-gradient(90deg,#c62828,#ef5350);}
    .reco-amber{background:linear-gradient(90deg,#e65100,#ff7043);}
    .reco-blue{background:linear-gradient(90deg,#1565c0,#42a5f5);}
    .reco-green{background:linear-gradient(90deg,#2e7d32,#66bb6a);}
</style>
""",
    unsafe_allow_html=True,
)

# ─────────────────────────────────────────────────────────────────────────────
# CONSTANTS
# ─────────────────────────────────────────────────────────────────────────────
STEPS: tuple[tuple[int, str, str], ...] = (
    (1, "📥", "Data Ingestion"),
    (2, "🔍", "Exploratory Analysis"),
    (3, "🧹", "Data Cleaning"),
    (4, "⚙️", "Feature Engineering"),
    (5, "🤖", "Model Training"),
    (6, "📊", "Model Evaluation"),
    (7, "📉", "Bias Detection"),
    (8, "💡", "Recommendations"),
    (9, "🎯", "Data Story"),
    (10, "📈", "Market Intelligence"),
)

# Zambia locale — currency (ZMW / Kwacha) and ID conventions
CURRENCY_CODE = "ZMW"
CURRENCY_LABEL = "Kwacha"
BILL_SHOCK_KWACHA = 600          # monthly bill shock threshold (K)
HIGH_VALUE_KWACHA = 400          # high-ARPU segment threshold (K)
LOYAL_VALUE_KWACHA = 350         # loyal / high-value usage threshold (K)
SAVE_RAR_KWACHA = 250            # revenue-at-risk save-list threshold (K)

# Zambia mobile operators — portfolio focus on Zamtel benchmarking
ZAMBIA_OPERATORS = ("Zamtel", "MTN Zambia", "Airtel Zambia")
FOCUS_OPERATOR = "Zamtel"

# Illustrative national market reference KPIs (Zambia mobile sector — for gap analysis)
ZAMBIA_MARKET_REFERENCE = pd.DataFrame({
    "operator": list(ZAMBIA_OPERATORS),
    "subscribers_m": [4.8, 9.2, 8.0],
    "market_share_pct": [22.0, 42.0, 36.0],
    "monthly_arpu_kwacha": [175.0, 215.0, 198.0],
    "churn_rate_pct": [3.4, 2.6, 2.9],
    "nps": [30.0, 41.0, 36.0],
    "csat": [3.8, 4.2, 4.0],
    "network_availability_pct": [93.5, 97.2, 96.0],
    "avg_latency_ms": [125.0, 95.0, 105.0],
})


def _fmt_kwacha(amount, decimals: int = 0) -> str:
    """Format an amount in Zambian Kwacha (K)."""
    if amount is None:
        return "N/A"
    try:
        val = float(amount)
        if np.isnan(val):
            return "N/A"
    except (TypeError, ValueError):
        return "N/A"
    if decimals == 0:
        return f"K {val:,.0f}"
    return f"K {val:,.2f}"


def _zambia_nrc(rng: np.random.Generator) -> str:
    """Zambian NRC format: NNNNNN/NN/N (National Registration Card)."""
    return f"{rng.integers(100000, 999999):06d}/{rng.integers(10, 99):02d}/{rng.integers(1, 9)}"


def _zambia_msisdn(rng: np.random.Generator, operator: str | None = None) -> str:
    """Zambian mobile in E.164: 260 + 9-digit (operator-specific prefixes)."""
    prefix_map = {
        "Zamtel": ["95", "96"],
        "MTN Zambia": ["76", "96"],
        "Airtel Zambia": ["77", "97"],
    }
    if operator and operator in prefix_map:
        prefix = rng.choice(prefix_map[operator])
    else:
        prefix = rng.choice(["95", "96", "76", "77", "97"])
    return f"260{prefix}{rng.integers(1000000, 9999999):07d}"


FIELD_DICT: tuple[tuple[str, str, str, str], ...] = (
    ("subscriber_id", "Identity", "Text", "Unique subscriber key  e.g. ZMB-SUB00001"),
    ("account_number", "Identity", "Text", "Billing account  e.g. ZMB-ACC123456"),
    ("msisdn", "Identity", "Text", "Mobile number E.164  e.g. 260971234567"),
    ("id_number", "Identity", "Text", "Zambian NRC  format NNNNNN/NN/N"),
    ("operator", "Market", "Text", "Mobile network operator — Zamtel / MTN Zambia / Airtel Zambia"),
    ("report_month", "Market", "Text", "Reporting period YYYY-MM for KPI trend monitoring"),
    ("plan", "Subscription", "Text", "Price plan name"),
    ("contract_type", "Subscription", "Text", "Month-to-Month / One Year / Two Year"),
    ("network_type", "Subscription", "Text", "2G / 3G / 4G LTE / 5G / NB-IoT"),
    ("service_status", "Subscription", "Text", "Active / Suspended / Terminated / Porting Out"),
    ("acquisition_channel", "Subscription", "Text", "How subscriber was acquired"),
    ("payment_method", "Subscription", "Text", "Debit Order / EFT / MTN Mobile Money / Airtel Money / Voucher"),
    ("region", "Subscription", "Text", "Province"),
    ("device_type", "Subscription", "Text", "Smartphone / Feature Phone / IoT Device etc."),
    ("node_id", "Network", "Text", "Base transceiver station ID"),
    ("tenure_months", "Tenure", "Integer", "Months since activation"),
    ("activation_date", "Tenure", "Date", "Service activation date  YYYY-MM-DD"),
    ("created_date", "Dates", "DateTime", "Ticket / transaction created  ISO 8601"),
    ("closed_date", "Dates", "DateTime", "Ticket / transaction closed  ISO 8601"),
    ("last_recharge_days_ago", "Tenure", "Integer", "Days since last recharge event"),
    ("sla_hours", "SLA", "Integer", "SLA target in hours  4/8/24/48"),
    ("monthly_charges", "Billing", "Float", f"Monthly invoice amount ({CURRENCY_CODE}) — may contain '-' tokens"),
    ("data_volume_mb", "Usage", "Float", "Total data consumed in MB — ~4% nulls injected"),
    ("data_overage_mb", "Usage", "Float", "Data consumed above plan bundle"),
    ("call_minutes", "Usage", "Float", "Voice call minutes — ~4% nulls injected"),
    ("sms_count", "Usage", "Integer", "SMS messages sent"),
    ("intl_call_minutes", "Usage", "Float", "International call minutes (0 if not roaming)"),
    ("roaming_flag", "Usage", "Binary", "1 = subscriber roamed in the period"),
    ("avg_recharge_amount", "Usage", "Float", "Average top-up amount (0 for postpaid)"),
    ("latency_ms", "Network KPI", "Float", "Average round-trip latency (ms)"),
    ("signal_strength_dbm", "Network KPI", "Float", "Average signal strength (dBm, -110 to -50)"),
    ("complaint_count", "CX", "Integer", "Complaints logged in period"),
    ("support_calls", "CX", "Integer", "Calls to support centre"),
    ("nps_score", "CX", "Integer", "Net Promoter Score  0-10"),
    ("csat_score", "CX", "Float", "Customer Satisfaction score  1.0-5.0"),
    ("app_logins_monthly", "CX", "Integer", "Self-service app logins per month"),
    ("refund_amount", "Billing", "Float", f"Refunds issued ({CURRENCY_CODE})"),
    ("payment_failures", "Billing", "Integer", "Failed payment attempts"),
    ("days_overdue", "Billing", "Integer", "Days payment is overdue"),
    ("bill_shock_flag", "Billing", "Binary", f"1 = bill > K{BILL_SHOCK_KWACHA} or data overage > 1 GB"),
    ("paperless_billing", "Billing", "Binary", "1 = enrolled in e-billing"),
    ("plan_upgrades", "Loyalty", "Integer", "Number of plan upgrades"),
    ("reactivation_count", "Loyalty", "Integer", "Times reactivated after suspension"),
    ("churn", "TARGET", "Binary", "1 = churned in period  (predict this)"),
)

PIPELINE_CHECKLIST: tuple[tuple[str, str, str], ...] = (
    ("Step 1 Ingestion", "All 44 columns", "CSV upload / demo data load"),
    ("Step 1 Ingestion", "operator, report_month", "Zambia operator tagging + KPI period"),
    ("Step 2 EDA", "All numeric columns", "Distributions, correlation heatmap, class balance"),
    ("Step 2 EDA", "tenure_months, acquisition_channel", "Cohort churn curves + channel quality"),
    ("Step 2 EDA", "monthly_charges", "Revenue-at-risk preview by tenure cohort"),
    ("Step 2 EDA", "data_volume_mb, latency_ms", "Missing value map (4% injected)"),
    ("Step 3 Cleaning", "monthly_charges", "Auto-coerce '-' tokens to numeric"),
    ("Step 3 Cleaning", "N/A, null, None tokens", "Null-token normalisation (payment_method, contract_type, service_status)"),
    ("Step 3 Cleaning", "Blank strings", "Whitespace strip + blank -> null (plan, region, device_type, channel)"),
    ("Step 3 Cleaning", "10 duplicate rows", "Duplicate row removal"),
    ("Step 3 Cleaning", "6 numeric cols", "Median/mean imputation of NaN values"),
    ("Step 4 Engineering", "monthly_charges, data_volume_mb", "Derived: charge_per_mb"),
    ("Step 4 Engineering", "complaint_count, tenure_months", "Derived: complaint_rate"),
    ("Step 4 Engineering", "call_minutes, sms_count, data_volume_mb", "Derived: engagement_score"),
    ("Step 4 Engineering", "monthly_charges, tenure_months", "Derived: estimated_clv, RFM-style scores"),
    ("Step 4 Engineering", "payment_failures, bill_shock_flag", "Secondary target: payment_risk"),
    ("Step 4 Engineering", "plan_upgrades, data_volume_mb", "Secondary target: upgrade_propensity"),
    ("Step 4 Engineering", "created_date, closed_date, sla_hours", "SLA status + resolution_hours"),
    ("Step 4 Engineering", "plan, region, device_type etc.", "One-hot encoding of categoricals"),
    ("Step 5 Training", "churn (target)", "Six models: LR, RF, GB, SVM, KNN, AdaBoost"),
    ("Step 5 Training", "All numeric features", "Feature importance (Random Forest)"),
    ("Step 6 Evaluation", "churn", "Confusion matrix, ROC-AUC, F1, precision, recall"),
    ("Step 6 Evaluation", "churn_probability", "Probability distribution histogram"),
    ("Step 7 Bias Detection", "All features", "Learning curves: overfit / underfit / good fit"),
    ("Step 7 Bias Detection", "Train vs Test", "Train/Test AUC gap per model"),
    ("Step 5 Training", "payment_risk, upgrade targets", "Secondary RF models: billing risk + upsell propensity"),
    ("Step 8 Recommendations", "revenue_at_risk, estimated_clv", "Save / Grow / Fix action segments + RAR dashboard"),
    ("Step 9 Data Story", "All stages", "Executive + technical PowerPoint decks + JSON export"),
    ("Step 10 Market Intelligence", "operator, monthly_charges, churn", "Zamtel vs MTN/Airtel KPI monitor + national benchmark gap"),
)

_STATE_DEFAULTS: dict = {
    "current_step": 1,
    "raw_df": None,
    "cleaned_df": None,
    "featured_df": None,
    "target_col": "churn",
    "feature_cols": [],
    "X_train": None,
    "X_test": None,
    "y_train": None,
    "y_test": None,
    "y_train_payment": None,
    "y_test_payment": None,
    "y_train_upgrade": None,
    "y_test_upgrade": None,
    "models": {},
    "model_results": {},
    "feature_importance": None,
    "recommendations_df": None,
    "secondary_models": {},
    "secondary_results": {},
    "analytics_kpis": {},
    "market_kpis": {},
    "focus_operator": FOCUS_OPERATOR,
    "imputer": None,
    "scaler": None,
    "steps_done": set(),
    "eda_insights": [],
    "clean_log": [],
    "dataset_name": "No dataset loaded",
    "demo": False,
}

SAMPLE_CSV_PATH = Path(__file__).resolve().parent / "telco_sample_data.csv"
LOGO_PATH = Path(__file__).resolve().parent / "analytics_logo.svg"

# ─────────────────────────────────────────────────────────────────────────────
# SESSION STATE
# ─────────────────────────────────────────────────────────────────────────────
ss = st.session_state
for _k, _v in _STATE_DEFAULTS.items():
    if _k not in ss:
        ss[_k] = _v if _k != "steps_done" else set()

# ─────────────────────────────────────────────────────────────────────────────
# UI HELPERS
# ─────────────────────────────────────────────────────────────────────────────
def _sidebar_logo() -> None:
    """Render Analytics logo in sidebar (HTML embed avoids Streamlit SVG file decode issues)."""
    if not LOGO_PATH.exists():
        return
    svg = LOGO_PATH.read_text(encoding="utf-8")
    st.sidebar.markdown(
        f'<div style="margin:0 0 12px 0;">{svg}</div>',
        unsafe_allow_html=True,
    )


def _nav_sidebar() -> None:
    _sidebar_logo()
    st.sidebar.markdown("## 📡 Zambia Telco Analytics")
    st.sidebar.caption(f"Zamtel benchmarking · {CURRENCY_CODE} · KPI monitor")

    done = ss.steps_done
    for num, icon, label in STEPS:
        unlocked = num == 1 or (num - 1) in done
        marker = "✅" if num in done else ("🔵" if num == ss.current_step else "○")
        if st.sidebar.button(
            f"{marker} {icon} {label}",
            key=f"nav_{num}",
            use_container_width=True,
            disabled=not unlocked,
        ):
            ss.current_step = num

    pct = int(len(done) / len(STEPS) * 100)
    st.sidebar.progress(pct / 100, text=f"Progress: {pct}%")

    if ss.raw_df is not None:
        st.sidebar.markdown("---")
        st.sidebar.markdown(
            f"**{ss.dataset_name}**  \n"
            f"{ss.raw_df.shape[0]:,} rows × {ss.raw_df.shape[1]} cols  \n"
            f"🎯 Target: `{ss.target_col}`"
        )

    if st.sidebar.button("🔄 Reset Pipeline", use_container_width=True):
        for k, v in _STATE_DEFAULTS.items():
            ss[k] = v if k != "steps_done" else set()
        st.rerun()


def _step_header(num: int, title: str, subtitle: str = "") -> None:
    icon = STEPS[num - 1][1]
    st.markdown(f'<div class="step-header">Step {num} {icon} {title}</div>', unsafe_allow_html=True)
    if subtitle:
        st.markdown(subtitle)


def _complete_step(num: int) -> None:
    ss.steps_done.add(num)
    ss.current_step = min(num + 1, len(STEPS))


def _kpi(label: str, value, color: str = "") -> None:
    cls = f"kpi-card {color}".strip()
    st.markdown(
        f'<div class="{cls}"><div class="kpi-value">{value}</div>'
        f'<div class="kpi-label">{label}</div></div>',
        unsafe_allow_html=True,
    )


def _insight(msg: str) -> None:
    st.markdown(f'<div class="insight-box">💡 {msg}</div>', unsafe_allow_html=True)


def _warn(msg: str) -> None:
    st.markdown(f'<div class="warn-box">⚠️ {msg}</div>', unsafe_allow_html=True)


def _ok(msg: str) -> None:
    st.markdown(f'<div class="ok-box">✔ {msg}</div>', unsafe_allow_html=True)


# Columns that should never be one-hot / label encoded as model inputs
_SKIP_ENCODE_COLS = {
    "subscriber_id", "msisdn", "id_number", "account_number",
    "activation_date", "created_date", "closed_date",
}
_MAX_OHE_CARDINALITY = 30


def _prepare_binary_target(target: pd.Series, col_name: str) -> pd.Series | None:
    """Coerce target to binary 0/1; return None if not usable for classification."""
    t = target.copy()
    if t.dtype == object:
        t = t.astype(str).str.strip()
        numeric = pd.to_numeric(t, errors="coerce")
        if numeric.notna().mean() > 0.9:
            t = numeric
        else:
            le = LabelEncoder()
            t = pd.Series(le.fit_transform(t), index=target.index, name=col_name)

    t = pd.to_numeric(t, errors="coerce")
    if t.isnull().any():
        _warn(f"Dropping {int(t.isnull().sum())} rows with missing target values.")

    unique = sorted(t.dropna().unique())
    if len(unique) == 1:
        st.error(f"Target `{col_name}` has only one class ({unique[0]}). Cannot train a classifier.")
        return None
    if len(unique) > 2:
        st.error(
            f"Target `{col_name}` has {len(unique)} classes ({unique[:8]}…). "
            "Select a binary churn column (0/1 or Yes/No)."
        )
        return None

    mapping = {unique[0]: 0, unique[1]: 1}
    return t.map(mapping).astype(int).rename(col_name)


def _safe_train_test_split(feat_df: pd.DataFrame, target: pd.Series):
    """Split with stratification when every class has at least 2 samples."""
    split_kwargs: dict = {"test_size": 0.2, "random_state": 42}
    counts = target.value_counts()
    if len(counts) >= 2 and counts.min() >= 2:
        split_kwargs["stratify"] = target
    else:
        rare = counts[counts < 2]
        _warn(
            "Stratified split skipped — some target classes have fewer than 2 samples "
            f"({', '.join(f'{k}: {v}' for k, v in rare.items())}). "
            "Using random split instead."
        )
    return train_test_split(feat_df, target, **split_kwargs)


def _derive_business_kpis(df: pd.DataFrame) -> pd.DataFrame:
    """Add CLV, RFM-style scores, and secondary prediction targets."""
    out = df.copy()

    if "monthly_charges" in out.columns and "tenure_months" in out.columns:
        out["estimated_clv"] = (out["monthly_charges"] * out["tenure_months"]).round(2)

    if "last_recharge_days_ago" in out.columns:
        out["rfm_recency"] = (90 - out["last_recharge_days_ago"].clip(0, 90)) / 90
    if "monthly_charges" in out.columns:
        charges = out["monthly_charges"]
        out["rfm_monetary"] = (charges - charges.min()) / (charges.max() - charges.min() + 1e-9)
    if "data_volume_mb" in out.columns:
        usage = out["data_volume_mb"]
        out["rfm_frequency"] = (usage - usage.min()) / (usage.max() - usage.min() + 1e-9)

    pay_fail = out.get("payment_failures", pd.Series(0, index=out.index)).fillna(0)
    overdue = out.get("days_overdue", pd.Series(0, index=out.index)).fillna(0)
    shock = out.get("bill_shock_flag", pd.Series(0, index=out.index)).fillna(0)
    out["_target_payment_risk"] = ((pay_fail > 0) | (overdue > 0) | (shock == 1)).astype(int)

    upgrades = out.get("plan_upgrades", pd.Series(0, index=out.index)).fillna(0)
    data_mb = out.get("data_volume_mb", pd.Series(0, index=out.index)).fillna(0)
    charges = out.get("monthly_charges", pd.Series(0, index=out.index)).fillna(0)
    churn_col = out.get("churn", pd.Series(0, index=out.index))
    churn_num = pd.to_numeric(churn_col, errors="coerce").fillna(0)
    out["_target_upgrade"] = (
        (upgrades > 0) | ((data_mb > 2500) & (charges < LOYAL_VALUE_KWACHA) & (churn_num == 0))
    ).astype(int)

    return out


def _apply_feature_pipeline(feat_df: pd.DataFrame) -> pd.DataFrame:
    """Apply stored imputer/scaler from Step 4 to a feature matrix."""
    X = feat_df.copy()
    if ss.imputer is not None:
        cols = X.columns.tolist()
        X = pd.DataFrame(ss.imputer.transform(X), columns=cols, index=X.index)
    if ss.scaler is not None:
        cols = X.columns.tolist()
        X = pd.DataFrame(ss.scaler.transform(X), columns=cols, index=X.index)
    return X


def _train_secondary_rf(
    X_train: pd.DataFrame,
    X_test: pd.DataFrame,
    y_train: pd.Series,
    y_test: pd.Series,
    label: str,
) -> tuple[RandomForestClassifier | None, dict | None]:
    """Train a Random Forest for a secondary binary target."""
    y_train = _prepare_binary_target(y_train, label)
    y_test = _prepare_binary_target(y_test, label)
    if y_train is None or y_test is None:
        return None, None
    if y_train.nunique() < 2:
        return None, None

    model = RandomForestClassifier(n_estimators=150, random_state=42, class_weight="balanced", n_jobs=-1)
    model.fit(X_train, y_train)
    y_prob = model.predict_proba(X_test)[:, 1]
    y_pred = model.predict(X_test)
    metrics = {
        "accuracy": accuracy_score(y_test, y_pred),
        "precision": precision_score(y_test, y_pred, zero_division=0),
        "recall": recall_score(y_test, y_pred, zero_division=0),
        "f1": f1_score(y_test, y_pred, zero_division=0),
        "roc_auc": roc_auc_score(y_test, y_prob),
        "positive_rate": float(y_train.mean()),
    }
    return model, metrics


def _action_segment(row) -> str:
    """Map scores to Save / Grow / Fix / Monitor actions."""
    churn_p = row.get("churn_probability", 0)
    rar = row.get("revenue_at_risk", 0)
    pay_p = row.get("payment_risk_probability", 0)
    up_p = row.get("upgrade_probability", 0)

    if churn_p >= 0.5 and rar >= SAVE_RAR_KWACHA:
        return "Save — High Value at Risk"
    if pay_p >= 0.45:
        return "Fix — Billing & Payment Risk"
    if up_p >= 0.45 and churn_p < 0.35:
        return "Grow — Upsell Ready"
    if churn_p >= 0.4:
        return "Watch — Elevated Churn Risk"
    return "Monitor — Stable"


def _operator_col(df: pd.DataFrame) -> str | None:
    """Return operator column name if present."""
    if "operator" in df.columns:
        return "operator"
    return None


def _compute_operator_kpis(df: pd.DataFrame, target_col: str = "churn") -> pd.DataFrame:
    """Aggregate portfolio KPIs by Zambian mobile operator."""
    op_col = _operator_col(df)
    if op_col is None:
        return pd.DataFrame()

    work = df.copy()
    work["_charges"] = pd.to_numeric(work.get("monthly_charges"), errors="coerce")
    work["_churn"] = pd.to_numeric(work.get(target_col), errors="coerce")
    work["_nps"] = pd.to_numeric(work.get("nps_score"), errors="coerce")
    work["_csat"] = pd.to_numeric(work.get("csat_score"), errors="coerce")
    work["_latency"] = pd.to_numeric(work.get("latency_ms"), errors="coerce")

    g = work.groupby(op_col, observed=True)
    out = g.size().reset_index(name="subscribers")
    out = out.rename(columns={op_col: "operator"})
    out["monthly_revenue"] = g["_charges"].sum().values
    out["avg_arpu"] = g["_charges"].mean().values
    if target_col in work.columns:
        out["churn_rate"] = g["_churn"].mean().values
        out["churn_rate_pct"] = (out["churn_rate"] * 100).round(2)
    if "nps_score" in work.columns:
        out["avg_nps"] = g["_nps"].mean().values
    if "csat_score" in work.columns:
        out["avg_csat"] = g["_csat"].mean().values
    if "latency_ms" in work.columns:
        out["avg_latency_ms"] = g["_latency"].mean().values
    out["market_share_pct"] = (out["subscribers"] / out["subscribers"].sum() * 100).round(1)
    return out


def _compute_monthly_trends(df: pd.DataFrame, target_col: str = "churn") -> pd.DataFrame:
    """KPI trends by report_month and operator."""
    if "report_month" not in df.columns or _operator_col(df) is None:
        return pd.DataFrame()

    work = df.copy()
    work["_charges"] = pd.to_numeric(work["monthly_charges"], errors="coerce")
    work["_churn"] = pd.to_numeric(work.get(target_col), errors="coerce")
    trend = (
        work.groupby(["report_month", "operator"], observed=True)
        .agg(
            subscribers=("operator", "count"),
            monthly_revenue=("_charges", "sum"),
            avg_arpu=("_charges", "mean"),
            churn_rate=("_churn", "mean"),
        )
        .reset_index()
    )
    trend["churn_rate_pct"] = (trend["churn_rate"] * 100).round(2)
    return trend.sort_values("report_month")


def _zamtel_gap_analysis(portfolio: pd.DataFrame) -> pd.DataFrame:
    """Compare Zamtel portfolio KPIs against national market reference."""
    if portfolio.empty or FOCUS_OPERATOR not in portfolio["operator"].values:
        return pd.DataFrame()

    ref = ZAMBIA_MARKET_REFERENCE.set_index("operator")
    zam = portfolio.set_index("operator").loc[FOCUS_OPERATOR]

    rows = []
    metrics = [
        ("monthly_arpu_kwacha", "avg_arpu", "ARPU (K)", True),
        ("churn_rate_pct", "churn_rate_pct", "Churn Rate (%)", False),
        ("nps", "avg_nps", "NPS", True),
        ("csat", "avg_csat", "CSAT", True),
        ("avg_latency_ms", "avg_latency_ms", "Latency (ms)", False),
        ("market_share_pct", "market_share_pct", "Market Share (%)", True),
    ]
    for ref_col, port_col, label, higher_better in metrics:
        if ref_col not in ref.columns or port_col not in zam:
            continue
        ref_val = float(ref.loc[FOCUS_OPERATOR, ref_col])
        port_val = float(zam[port_col])
        if np.isnan(port_val):
            continue
        gap = port_val - ref_val
        if higher_better:
            status = "Above benchmark" if gap >= 0 else "Below benchmark"
        else:
            status = "Above benchmark" if gap <= 0 else "Below benchmark"
        rows.append({
            "KPI": label,
            "Zamtel (portfolio)": round(port_val, 2),
            "National reference": round(ref_val, 2),
            "Gap": round(gap, 2),
            "Status": status,
        })
    return pd.DataFrame(rows)


def _render_market_kpi_cards(kpis: pd.DataFrame, focus: str) -> None:
    """Top-row KPI cards for a selected operator."""
    if kpis.empty or focus not in kpis["operator"].values:
        st.caption("Operator KPIs unavailable — ensure `operator` column is populated.")
        return

    row = kpis.loc[kpis["operator"] == focus].iloc[0]
    c1, c2, c3, c4, c5, c6 = st.columns(6)
    with c1:
        _kpi("Subscribers", f"{int(row['subscribers']):,}")
    with c2:
        rev = row.get("monthly_revenue", 0)
        _kpi("Monthly Revenue", _fmt_kwacha(rev), "green")
    with c3:
        _kpi("ARPU", _fmt_kwacha(row.get("avg_arpu", 0)), "purple")
    with c4:
        churn = row.get("churn_rate_pct", row.get("churn_rate", 0) * 100 if "churn_rate" in row else 0)
        _kpi("Churn Rate", f"{churn:.1f}%", "orange")
    with c5:
        nps = row.get("avg_nps", np.nan)
        _kpi("Avg NPS", f"{nps:.0f}" if pd.notna(nps) else "N/A")
    with c6:
        share = row.get("market_share_pct", 0)
        _kpi("Portfolio Share", f"{share:.1f}%", "green")


def _render_market_snapshot(df: pd.DataFrame, target_col: str, expanded: bool = False) -> None:
    """Compact Zambia market KPI panel (Steps 1–2)."""
    op_col = _operator_col(df)
    if op_col is None:
        with st.expander("🇿🇲 Zambia Market KPIs", expanded=expanded):
            _warn(
                "No `operator` column found. Upload data with Zamtel / MTN Zambia / Airtel Zambia "
                "or use **Demo Data** for a Zambia-localised sample."
            )
        return

    kpis = _compute_operator_kpis(df, target_col)
    ss.market_kpis = kpis.to_dict("records")

    with st.expander("🇿🇲 Zambia Market KPI Monitor", expanded=expanded):
        focus = st.selectbox(
            "Focus operator",
            list(kpis["operator"]),
            index=list(kpis["operator"]).index(FOCUS_OPERATOR)
            if FOCUS_OPERATOR in kpis["operator"].values else 0,
            key=f"market_focus_{expanded}",
        )
        ss.focus_operator = focus
        _render_market_kpi_cards(kpis, focus)

        col_a, col_b = st.columns(2)
        with col_a:
            fig_share = px.pie(
                kpis, names="operator", values="subscribers",
                title="Subscriber Base by Operator (portfolio)",
                color="operator",
                color_discrete_map={
                    "Zamtel": "#00695c",
                    "MTN Zambia": "#ffcc00",
                    "Airtel Zambia": "#e53935",
                },
            )
            st.plotly_chart(fig_share, use_container_width=True)
        with col_b:
            fig_arpu = px.bar(
                kpis, x="operator", y="avg_arpu", text_auto=".0f",
                title="Average ARPU by Operator (K)",
                color="operator",
                color_discrete_map={
                    "Zamtel": "#00695c",
                    "MTN Zambia": "#ffcc00",
                    "Airtel Zambia": "#e53935",
                },
            )
            st.plotly_chart(fig_arpu, use_container_width=True)

        if focus == FOCUS_OPERATOR:
            gap = _zamtel_gap_analysis(kpis)
            if not gap.empty:
                st.markdown(f"**{FOCUS_OPERATOR} vs National Market Reference**")
                st.dataframe(gap, use_container_width=True, hide_index=True)
                below = gap[gap["Status"] == "Below benchmark"]["KPI"].tolist()
                if below:
                    _warn(f"Gap areas for {FOCUS_OPERATOR}: **{', '.join(below)}**")
                else:
                    _ok(f"{FOCUS_OPERATOR} meets or exceeds all tracked national benchmarks in this portfolio.")


def _render_operator_benchmark_charts(df: pd.DataFrame, target_col: str) -> None:
    """Full operator comparison charts for Market Intelligence step."""
    kpis = _compute_operator_kpis(df, target_col)
    if kpis.empty:
        st.warning("Operator benchmarking requires an `operator` column.")
        return

    st.subheader("Operator Benchmark — Zamtel vs MTN vs Airtel")
    tab_port, tab_ref, tab_trend, tab_region = st.tabs(
        ["Portfolio Comparison", "National Reference", "KPI Trends", "Regional View"]
    )

    with tab_port:
        m1, m2 = st.columns(2)
        metrics_long = kpis.melt(
            id_vars=["operator"],
            value_vars=[c for c in ["avg_arpu", "churn_rate_pct", "avg_nps", "avg_csat"] if c in kpis.columns],
            var_name="metric",
            value_name="value",
        )
        with m1:
            fig_rev = px.bar(
                kpis.sort_values("monthly_revenue", ascending=False),
                x="operator", y="monthly_revenue", text_auto=".2s",
                title="Monthly Revenue by Operator (portfolio)",
                color="operator",
                color_discrete_map={"Zamtel": "#00695c", "MTN Zambia": "#ffcc00", "Airtel Zambia": "#e53935"},
            )
            st.plotly_chart(fig_rev, use_container_width=True)
        with m2:
            fig_churn = px.bar(
                kpis.sort_values("churn_rate_pct", ascending=True),
                x="operator", y="churn_rate_pct", text_auto=".2f",
                title="Churn Rate by Operator (%)",
                color="churn_rate_pct", color_continuous_scale="RdYlGn_r",
            )
            st.plotly_chart(fig_churn, use_container_width=True)

        fig_radar = px.line_polar(
            metrics_long, r="value", theta="metric", color="operator", line_close=True,
            title="Multi-KPI Operator Profile (portfolio normalised scale)",
            color_discrete_map={"Zamtel": "#00695c", "MTN Zambia": "#ffcc00", "Airtel Zambia": "#e53935"},
        )
        st.plotly_chart(fig_radar, use_container_width=True)

        gap = _zamtel_gap_analysis(kpis)
        if not gap.empty:
            st.markdown(f"### {FOCUS_OPERATOR} Gap Analysis (portfolio vs national reference)")
            fig_gap = px.bar(
                gap, x="KPI", y="Gap", color="Status",
                color_discrete_map={"Above benchmark": "#2e7d32", "Below benchmark": "#c62828"},
                title=f"{FOCUS_OPERATOR} KPI Gap vs National Benchmark",
            )
            st.plotly_chart(fig_gap, use_container_width=True)
            st.dataframe(gap, use_container_width=True, hide_index=True)

    with tab_ref:
        st.caption(
            "National reference figures are illustrative industry benchmarks for Zambia's mobile market "
            "(millions of subscribers at national scale). Compare directionally against your portfolio KPIs."
        )
        ref = ZAMBIA_MARKET_REFERENCE.copy()
        ref_display = ref.rename(columns={
            "subscribers_m": "Subscribers (M, national)",
            "monthly_arpu_kwacha": "ARPU (K)",
            "churn_rate_pct": "Churn (%)",
            "network_availability_pct": "Network Availability (%)",
        })
        st.dataframe(ref_display, use_container_width=True, hide_index=True)

        ref_melt = ref.melt(
            id_vars=["operator"],
            value_vars=["market_share_pct", "monthly_arpu_kwacha", "churn_rate_pct", "nps"],
            var_name="KPI", value_name="Value",
        )
        fig_ref = px.bar(
            ref_melt, x="KPI", y="Value", color="operator", barmode="group",
            title="National Market Reference by Operator",
            color_discrete_map={"Zamtel": "#00695c", "MTN Zambia": "#ffcc00", "Airtel Zambia": "#e53935"},
        )
        st.plotly_chart(fig_ref, use_container_width=True)

        if FOCUS_OPERATOR in kpis["operator"].values:
            zam_row = kpis.loc[kpis["operator"] == FOCUS_OPERATOR].iloc[0]
            ref_row = ref.loc[ref["operator"] == FOCUS_OPERATOR].iloc[0]
            comp = pd.DataFrame([
                {"Metric": "ARPU (K)", "Portfolio": zam_row.get("avg_arpu"), "National ref": ref_row["monthly_arpu_kwacha"]},
                {"Metric": "Churn (%)", "Portfolio": zam_row.get("churn_rate_pct"), "National ref": ref_row["churn_rate_pct"]},
                {"Metric": "NPS", "Portfolio": zam_row.get("avg_nps"), "National ref": ref_row["nps"]},
                {"Metric": "Share (%)", "Portfolio": zam_row.get("market_share_pct"), "National ref": ref_row["market_share_pct"]},
            ])
            fig_comp = px.bar(
                comp.melt(id_vars="Metric", var_name="Source", value_name="Value"),
                x="Metric", y="Value", color="Source", barmode="group",
                title=f"{FOCUS_OPERATOR}: Portfolio vs National Reference",
            )
            st.plotly_chart(fig_comp, use_container_width=True)

    with tab_trend:
        trend = _compute_monthly_trends(df, target_col)
        if trend.empty:
            st.info("Add a `report_month` column (YYYY-MM) to enable KPI trend monitoring.")
        else:
            t1, t2 = st.columns(2)
            with t1:
                fig_subs = px.line(
                    trend, x="report_month", y="subscribers", color="operator", markers=True,
                    title="Subscriber Base Trend by Operator",
                    color_discrete_map={"Zamtel": "#00695c", "MTN Zambia": "#ffcc00", "Airtel Zambia": "#e53935"},
                )
                st.plotly_chart(fig_subs, use_container_width=True)
            with t2:
                fig_rev_t = px.line(
                    trend, x="report_month", y="monthly_revenue", color="operator", markers=True,
                    title="Monthly Revenue Trend (K)",
                    color_discrete_map={"Zamtel": "#00695c", "MTN Zambia": "#ffcc00", "Airtel Zambia": "#e53935"},
                )
                st.plotly_chart(fig_rev_t, use_container_width=True)
            fig_churn_t = px.line(
                trend, x="report_month", y="churn_rate_pct", color="operator", markers=True,
                title="Churn Rate Trend (%)",
                color_discrete_map={"Zamtel": "#00695c", "MTN Zambia": "#ffcc00", "Airtel Zambia": "#e53935"},
            )
            st.plotly_chart(fig_churn_t, use_container_width=True)

    with tab_region:
        if "region" in df.columns and _operator_col(df):
            reg = df.copy()
            reg["_churn"] = pd.to_numeric(reg.get(target_col), errors="coerce")
            reg["_charges"] = pd.to_numeric(reg["monthly_charges"], errors="coerce")
            reg_kpi = (
                reg.groupby(["region", "operator"], observed=True)
                .agg(subscribers=("operator", "count"), revenue=("_charges", "sum"), churn=("_churn", "mean"))
                .reset_index()
            )
            fig_reg = px.bar(
                reg_kpi, x="region", y="subscribers", color="operator", barmode="group",
                title="Subscribers by Province & Operator",
                color_discrete_map={"Zamtel": "#00695c", "MTN Zambia": "#ffcc00", "Airtel Zambia": "#e53935"},
            )
            st.plotly_chart(fig_reg, use_container_width=True)
        else:
            st.info("Regional operator view requires `region` and `operator` columns.")


# ─────────────────────────────────────────────────────────────────────────────
# DEMO DATA
# ─────────────────────────────────────────────────────────────────────────────
def _demo_data(n: int = 2000) -> pd.DataFrame:
    """
    Full-featured synthetic Zambia telco dataset with 44 columns covering:
    Identity, Market, Subscription, Usage, Network KPIs, Customer Experience,
    Billing, Loyalty — plus injected nulls, bad tokens & duplicates
    to exercise every pipeline step.
    """
    rng = np.random.default_rng(42)
    operators = list(ZAMBIA_OPERATORS)
    op_weights = [0.34, 0.38, 0.28]  # portfolio skew toward Zamtel CRM extract
    report_months = ["2025-01", "2025-02", "2025-03", "2025-04", "2025-05", "2025-06"]
    plans = [
        "Prepaid Basic", "Prepaid Smart", "Postpaid K199", "Postpaid K499",
        "Postpaid K799", "IoT SIM", "Business Pro",
    ]
    regions = [
        "Lusaka", "Copperbelt", "Southern", "Eastern", "Northern",
        "North-Western", "Western", "Luapula", "Muchinga", "Central",
    ]
    devices = ["Smartphone", "Feature Phone", "IoT Device", "Tablet", "Router", "Wearable"]
    networks = ["4G LTE", "3G", "5G", "2G", "NB-IoT"]
    channels = ["Direct Sales", "Online Portal", "Retail Store", "Agent", "Referral"]
    contracts = ["Month-to-Month", "One Year", "Two Year"]
    payments = ["Debit Order", "EFT", "MTN Mobile Money", "Airtel Money", "Voucher"]
    statuses = ["Active", "Suspended", "Terminated", "Porting Out"]
    nodes = [
        "BTS_LSK_001", "BTS_NDL_007", "BTS_KTW_003",
        "BTS_LVN_012", "BTS_KAB_005", "BTS_CHP_009",
    ]

    idx = np.arange(n)
    base_date = datetime(2024, 1, 1)
    operator_assign = rng.choice(operators, n, p=op_weights)

    tenure = rng.integers(1, 72, n)
    activation = [base_date + timedelta(days=int(rng.integers(0, 900))) for _ in idx]
    created = [a + timedelta(hours=int(rng.integers(1, 500))) for a in activation]
    closed = [c + timedelta(hours=int(rng.integers(1, 96))) for c in created]

    monthly_charges = rng.uniform(79, 850, n).round(2)
    # Operator-specific portfolio characteristics for benchmarking
    for op, adj in [("Zamtel", -35), ("MTN Zambia", 25), ("Airtel Zambia", 5)]:
        mask = operator_assign == op
        monthly_charges[mask] = (monthly_charges[mask] + adj).clip(50, 900)
    data_volume = rng.uniform(50, 8000, n).round(1)
    call_minutes = rng.uniform(10, 500, n).round(1)
    complaint_count = rng.integers(0, 6, n)
    churn = (
        (complaint_count >= 3) | (monthly_charges > 650) | (data_volume < 200)
    ).astype(int)
    churn = np.where(rng.random(n) < 0.12, 1 - churn, churn)
    # Zamtel slightly higher churn in synthetic portfolio
    churn = np.where((operator_assign == "Zamtel") & (rng.random(n) < 0.06), 1, churn)

    nps_base = rng.integers(0, 11, n)
    latency_base = rng.uniform(80, 350, n).round(1)
    for op, nps_adj, lat_adj in [("Zamtel", -4, 25), ("MTN Zambia", 3, -15), ("Airtel Zambia", 0, 5)]:
        mask = operator_assign == op
        nps_base[mask] = np.clip(nps_base[mask] + nps_adj, 0, 10)
        latency_base[mask] = np.clip(latency_base[mask] + lat_adj, 60, 400)

    df = pd.DataFrame({
        "subscriber_id": [f"ZMB-SUB{i:05d}" for i in idx],
        "account_number": [f"ZMB-ACC{rng.integers(100000, 999999)}" for _ in idx],
        "msisdn": [_zambia_msisdn(rng, op) for op in operator_assign],
        "id_number": [_zambia_nrc(rng) for _ in idx],
        "operator": operator_assign,
        "report_month": rng.choice(report_months, n),
        "plan": rng.choice(plans, n),
        "contract_type": rng.choice(contracts, n),
        "network_type": rng.choice(networks, n),
        "service_status": rng.choice(statuses, n, p=[0.82, 0.08, 0.07, 0.03]),
        "acquisition_channel": rng.choice(channels, n),
        "payment_method": rng.choice(payments, n),
        "region": rng.choice(regions, n),
        "device_type": rng.choice(devices, n),
        "node_id": rng.choice(nodes, n),
        "tenure_months": tenure,
        "activation_date": [d.strftime("%Y-%m-%d") for d in activation],
        "created_date": [d.strftime("%Y-%m-%dT%H:%M:%S") for d in created],
        "closed_date": [d.strftime("%Y-%m-%dT%H:%M:%S") for d in closed],
        "last_recharge_days_ago": rng.integers(1, 90, n),
        "sla_hours": rng.choice([4, 8, 24, 48], n),
        "monthly_charges": monthly_charges,
        "data_volume_mb": data_volume,
        "data_overage_mb": rng.uniform(0, 500, n).round(1),
        "call_minutes": call_minutes,
        "sms_count": rng.integers(0, 80, n),
        "intl_call_minutes": np.where(rng.random(n) < 0.15, rng.uniform(0, 200, n).round(2), 0),
        "roaming_flag": rng.integers(0, 2, n),
        "avg_recharge_amount": np.where(
            rng.random(n) < 0.4, rng.uniform(10, 150, n).round(2), 0
        ),
        "latency_ms": latency_base,
        "signal_strength_dbm": rng.uniform(-110, -50, n).round(1),
        "complaint_count": complaint_count,
        "support_calls": rng.integers(0, 5, n),
        "nps_score": nps_base,
        "csat_score": rng.uniform(1, 5, n).round(1),
        "app_logins_monthly": rng.integers(0, 30, n),
        "refund_amount": np.where(rng.random(n) < 0.08, rng.uniform(10, 350, n).round(2), 0),
        "payment_failures": rng.integers(0, 4, n),
        "days_overdue": rng.integers(0, 45, n),
        "bill_shock_flag": (monthly_charges > BILL_SHOCK_KWACHA).astype(int),
        "paperless_billing": rng.integers(0, 2, n),
        "plan_upgrades": rng.integers(0, 3, n),
        "reactivation_count": rng.integers(0, 2, n),
        "churn": churn,
    })

    # Inject nulls (~4% on selected numerics)
    null_cols = ["data_volume_mb", "call_minutes", "refund_amount", "latency_ms", "csat_score", "signal_strength_dbm"]
    for col in null_cols:
        mask = rng.random(n) < 0.04
        df.loc[mask, col] = np.nan

    # Bad tokens (monthly_charges must be object dtype to hold '-')
    df["monthly_charges"] = df["monthly_charges"].astype(object)
    bad_idx = rng.choice(n, 40, replace=False)
    df.loc[bad_idx[:15], "monthly_charges"] = "-"
    df.loc[bad_idx[15:25], "payment_method"] = "N/A"
    df.loc[bad_idx[25:30], "contract_type"] = "null"
    df.loc[bad_idx[30:35], "service_status"] = "None"
    df.loc[bad_idx[35:], "plan"] = "  "

    # Duplicates
    dup = df.iloc[:10].copy()
    df = pd.concat([df, dup], ignore_index=True)

    return df


# ─────────────────────────────────────────────────────────────────────────────
# STEP 1 — DATA INGESTION
# ─────────────────────────────────────────────────────────────────────────────
def step_ingest() -> None:
    _step_header(
        1, "Data Ingestion",
        "Upload your Zambia telecom dataset (CSV / Excel), download the ready-made sample CSV to test "
        "all checklist items, or click **Use Demo Data** to load a synthetic 2,000-subscriber "
        "Zambia portfolio (Zamtel · MTN · Airtel) with KPI monitoring fields.",
    )

    with st.expander("📋  Full Field Dictionary — 44 Telco Columns", expanded=False):
        fd = pd.DataFrame(FIELD_DICT, columns=["Column", "Category", "Type", "Description"])
        st.dataframe(fd, use_container_width=True, hide_index=True, height=420)
        cat_summary = (
            fd.groupby("Category")["Column"].count().reset_index(name="Fields")
        )
        fig_cat = px.bar(
            cat_summary, x="Category", y="Fields", text="Fields",
            title="Fields by Category", color="Fields", color_continuous_scale="Blues",
        )
        st.plotly_chart(fig_cat, use_container_width=True)

    with st.expander("✅  Pipeline Checklist — What each field tests", expanded=False):
        cl = pd.DataFrame(PIPELINE_CHECKLIST, columns=["Step", "Field(s) Used", "What is Tested"])
        st.dataframe(cl, use_container_width=True, hide_index=True)

    col_up, col_demo, col_dl = st.columns([2, 1, 1])
    with col_up:
        uploaded = st.file_uploader(
            "Upload CSV or Excel file", type=["csv", "xlsx", "xls"], label_visibility="collapsed"
        )
    with col_demo:
        if st.button("Use Demo Data", type="primary", use_container_width=True):
            df = _demo_data(2000)
            ss.raw_df = df
            ss.dataset_name = "Zambia Telco Demo (2,000 rows × 44 cols · Zamtel focus)"
            ss.demo = True
            ss.target_col = "churn" if "churn" in df.columns else df.columns[-1]
            ss.cleaned_df = None
            ss.featured_df = None
            ss.model_results = {}
            ss.models = {}
            _ok(f"Demo dataset loaded: {df.shape[0]:,} rows × {df.shape[1]} cols")
    with col_dl:
        if SAMPLE_CSV_PATH.exists():
            st.download_button(
                "Download Sample CSV",
                data=SAMPLE_CSV_PATH.read_bytes(),
                file_name="telco_sample_data.csv",
                mime="text/csv",
                use_container_width=True,
                help="500-row CSV with all 44 fields — upload it back to test every step",
            )
        else:
            st.caption("Sample CSV not found")

    if uploaded is not None:
        try:
            name = uploaded.name.lower()
            if name.endswith(".csv"):
                df = pd.read_csv(uploaded)
            else:
                df = pd.read_excel(uploaded)
            ss.raw_df = df
            ss.dataset_name = uploaded.name
            ss.demo = False
            if "churn" in df.columns:
                ss.target_col = "churn"
            ss.cleaned_df = None
            ss.featured_df = None
            ss.model_results = {}
            ss.models = {}
            _ok(f"Loaded **{uploaded.name}**: {df.shape[0]:,} rows × {df.shape[1]} cols")
        except Exception as exc:
            st.error(f"Could not read file: {exc}")

    if ss.raw_df is None:
        st.info("Load data above to continue.")
        return

    df = ss.raw_df
    st.subheader("Data Preview")
    st.dataframe(df.head(20), use_container_width=True, height=320)

    c1, c2, c3, c4 = st.columns(4)
    with c1:
        _kpi("Total Records", f"{df.shape[0]:,}")
    with c2:
        _kpi("Columns", df.shape[1], "green")
    with c3:
        _kpi("Missing Cells", f"{int(df.isnull().sum().sum()):,}", "orange")
    with c4:
        _kpi("Duplicates", f"{int(df.duplicated().sum()):,}", "purple")

    _render_market_snapshot(df, ss.target_col, expanded=True)

    st.subheader("Column Inventory")
    inv = pd.DataFrame({
        "Column": df.columns,
        "Type": df.dtypes.astype(str).values,
        "Non-Null": df.notnull().sum().values,
        "Null %": (df.isnull().mean() * 100).round(1).values,
        "Unique": df.nunique().values,
        "Sample": [str(df[c].dropna().iloc[0]) if df[c].notnull().any() else "" for c in df.columns],
    })
    st.dataframe(inv, use_container_width=True, hide_index=True)

    st.subheader("Target Column (for Churn Prediction)")
    tgt_options = list(df.columns)
    default_idx = tgt_options.index(ss.target_col) if ss.target_col in tgt_options else 0
    ss.target_col = st.selectbox(
        "Select the churn / target column (must be 0/1 or binary):",
        tgt_options,
        index=default_idx,
    )

    target_series = df[ss.target_col].dropna()
    if target_series.dtype == object:
        unique_vals = sorted(target_series.unique())
        _warn(
            f"Target `{ss.target_col}` is text. Values: {unique_vals[:5]}. "
            "It will be label-encoded in the Feature Engineering step."
        )
    elif target_series.nunique() > 2:
        _warn(
            f"Target `{ss.target_col}` has {target_series.nunique()} distinct values "
            f"({sorted(target_series.unique())[:8]}…). "
            "Churn prediction requires a **binary** column (0/1)."
        )

    if st.button("Confirm Ingestion & Proceed →", type="primary"):
        _complete_step(1)
        st.rerun()





# ─────────────────────────────────────────────────────────────────────────────
# STEP 2 — EXPLORATORY DATA ANALYSIS
# ─────────────────────────────────────────────────────────────────────────────
def step_eda() -> None:
    _step_header(
        2, "Exploratory Data Analysis",
        "Understand distributions, correlations, missing values, and target balance before cleaning.",
    )
    if ss.raw_df is None:
        st.warning("Complete Step 1 first.")
        return

    df = ss.raw_df.copy()
    tgt = ss.target_col
    ss.eda_insights = []

    st.subheader("Descriptive Statistics")
    num_cols = df.select_dtypes(include="number").columns.tolist()
    if num_cols:
        st.dataframe(df[num_cols].describe().T.round(3), use_container_width=True)
    else:
        st.caption("No numeric columns detected.")

    st.subheader(f"Target Distribution — {tgt}")
    if tgt in df.columns:
        vc = df[tgt].value_counts(dropna=False).reset_index()
        vc.columns = [tgt, "count"]
        vc["pct"] = (vc["count"] / vc["count"].sum() * 100).round(1).astype(str) + "%"
        col_pie, col_bar = st.columns(2)
        with col_pie:
            fig1 = px.pie(vc, names=tgt, values="count", title="Class Proportions")
            st.plotly_chart(fig1, use_container_width=True)
        with col_bar:
            bar_color = tgt if vc.shape[0] < 10 else None
            fig2 = px.bar(vc, x=tgt, y="count", text="pct", color=bar_color, title="Class Counts")
            st.plotly_chart(fig2, use_container_width=True)
        churn_rate = df[tgt].astype(float).mean() if pd.to_numeric(df[tgt], errors="coerce").notna().any() else None
        if churn_rate is not None:
            _insight(f"Overall churn / positive-class rate: **{churn_rate:.1%}**")
            ss.eda_insights.append(f"Churn rate: {churn_rate:.1%}")
    else:
        st.error(f"Target column `{tgt}` not found.")

    st.subheader("Correlation Heatmap")
    if len(num_cols) >= 2:
        corr = df[num_cols].corr()
        fig3 = px.imshow(
            corr, text_auto=".2f", aspect="auto",
            color_continuous_scale="RdBu_r", zmin=-1, zmax=1,
            title="Pearson Correlation Matrix",
        )
        st.plotly_chart(fig3, use_container_width=True)
        if tgt in num_cols:
            top_corr = corr[tgt].drop(tgt, errors="ignore").abs().sort_values(ascending=False).head(5)
            if len(top_corr):
                pairs = ", ".join(f"`{k}` ({v:.2f})" for k, v in top_corr.items())
                _insight(f"Top correlations with `{tgt}`: {pairs}")
    else:
        st.caption("Need at least 2 numeric columns for correlation.")

    st.subheader("Missing Values")
    miss = df.isnull().sum()
    miss = miss[miss > 0].sort_values(ascending=False)
    if len(miss):
        miss_df = miss.reset_index()
        miss_df.columns = ["Column", "Missing"]
        miss_df["Pct"] = (miss_df["Missing"] / len(df) * 100).round(1)
        fig4 = px.bar(miss_df, x="Column", y="Missing", color="Pct",
                      color_continuous_scale="Oranges", title="Missing Value Counts")
        st.plotly_chart(fig4, use_container_width=True)
        _warn(f"{len(miss)} columns contain null values — address in Step 3.")
        ss.eda_insights.append(f"{len(miss)} columns with missing values")
    else:
        _ok("No missing values detected in raw data.")

    st.subheader("Numeric Distributions")
    if num_cols:
        dist_col = st.selectbox("Numeric column for distribution:", num_cols, key="eda_dist_col")
        if dist_col:
            fig_dist = px.histogram(
                df, x=dist_col, color=tgt if tgt in df.columns else None,
                nbins=40, marginal="box", title=f"Distribution — {dist_col}",
            )
            st.plotly_chart(fig_dist, use_container_width=True)
            skew = df[dist_col].skew()
            _insight(f"`{dist_col}` skewness: **{skew:.2f}**" + (" (right-skewed)" if skew > 1 else ""))

    st.subheader("Target Correlation Ranking")
    if tgt in num_cols:
        tgt_corr = (
            df[num_cols].corr()[tgt].drop(tgt, errors="ignore").abs().sort_values(ascending=False)
        )
        corr_df = tgt_corr.head(12).reset_index()
        corr_df.columns = ["Feature", "|Correlation|"]
        fig_tc = px.bar(corr_df, x="|Correlation|", y="Feature", orientation="h",
                        title=f"Top Features Correlated with {tgt}")
        st.plotly_chart(fig_tc, use_container_width=True)

    st.subheader("Categorical Breakdown")
    cat_cols = [c for c in df.select_dtypes("object").columns if c != tgt][:4]
    if cat_cols and tgt in df.columns:
        pick = st.selectbox("Categorical column to explore:", cat_cols)
        if pick:
            ct = df.groupby(pick, dropna=False)[tgt].apply(
                lambda s: pd.to_numeric(s, errors="coerce").mean()
            ).reset_index()
            ct.columns = [pick, "churn_rate"]
            fig5 = px.bar(ct.sort_values("churn_rate", ascending=False),
                          x=pick, y="churn_rate", title=f"Churn Rate by {pick}")
            st.plotly_chart(fig5, use_container_width=True)

    st.subheader("Tenure Cohort Analysis")
    if "tenure_months" in df.columns and tgt in df.columns:
        cohort_df = df.copy()
        cohort_df["_tenure_bin"] = pd.cut(
            cohort_df["tenure_months"],
            bins=[0, 6, 12, 24, 48, 999],
            labels=["0-6m", "6-12m", "12-24m", "24-48m", "48m+"],
        )
        cohort_churn = (
            cohort_df.groupby("_tenure_bin", observed=True)[tgt]
            .apply(lambda s: pd.to_numeric(s, errors="coerce").mean())
            .reset_index()
        )
        cohort_churn.columns = ["Tenure Cohort", "Churn Rate"]
        fig_cohort = px.bar(
            cohort_churn, x="Tenure Cohort", y="Churn Rate",
            text=cohort_churn["Churn Rate"].map(lambda x: f"{x:.1%}"),
            color="Churn Rate", color_continuous_scale="Reds",
            title="Churn Rate by Tenure Cohort",
        )
        fig_cohort.update_traces(textposition="outside")
        st.plotly_chart(fig_cohort, use_container_width=True)
        worst = cohort_churn.loc[cohort_churn["Churn Rate"].idxmax(), "Tenure Cohort"]
        _insight(f"Highest churn cohort: **{worst}** — prioritise onboarding/early-life retention.")
        ss.eda_insights.append(f"Highest churn tenure cohort: {worst}")

        if "monthly_charges" in df.columns:
            cohort_df["charges"] = pd.to_numeric(cohort_df["monthly_charges"], errors="coerce")
            cohort_df["churn_num"] = pd.to_numeric(cohort_df[tgt], errors="coerce")
            cohort_df["cohort_rar"] = cohort_df["charges"] * cohort_df["churn_num"]
            rar_by_cohort = (
                cohort_df.groupby("_tenure_bin", observed=True)["cohort_rar"]
                .sum().reset_index()
            )
            rar_by_cohort.columns = ["Tenure Cohort", "Revenue Lost (actual)"]
            fig_rar = px.bar(
                rar_by_cohort, x="Tenure Cohort", y="Revenue Lost (actual)",
                title="Actual Revenue Lost by Tenure Cohort (charges × churned)",
                color="Revenue Lost (actual)", color_continuous_scale="Oranges",
            )
            st.plotly_chart(fig_rar, use_container_width=True)
            total_lost = cohort_df["cohort_rar"].sum()
            _insight(f"Observed revenue lost to churn in sample: **{_fmt_kwacha(total_lost)}**")
            ss.analytics_kpis["observed_revenue_lost"] = float(total_lost)

    st.subheader("Acquisition Channel Quality")
    if "acquisition_channel" in df.columns and tgt in df.columns:
        ch = df.groupby("acquisition_channel")[tgt].apply(
            lambda s: pd.to_numeric(s, errors="coerce").mean()
        ).reset_index()
        ch.columns = ["Channel", "Churn Rate"]
        ch = ch.sort_values("Churn Rate", ascending=False)
        fig_ch = px.bar(
            ch, x="Channel", y="Churn Rate",
            color="Churn Rate", color_continuous_scale="RdYlGn_r",
            title="Churn Rate by Acquisition Channel",
        )
        st.plotly_chart(fig_ch, use_container_width=True)
        best_ch = ch.loc[ch["Churn Rate"].idxmin(), "Channel"]
        worst_ch = ch.loc[ch["Churn Rate"].idxmax(), "Channel"]
        _insight(f"Best channel: **{best_ch}** · Worst channel: **{worst_ch}**")
        ss.eda_insights.append(f"Best acquisition channel: {best_ch}; worst: {worst_ch}")

    if _operator_col(df):
        st.subheader("Zambia Operator Snapshot")
        op_kpis = _compute_operator_kpis(df, tgt)
        oc1, oc2 = st.columns(2)
        with oc1:
            fig_op_churn = px.bar(
                op_kpis, x="operator", y="churn_rate_pct", text_auto=".1f",
                title="Churn Rate by Operator (%)",
                color="operator",
                color_discrete_map={"Zamtel": "#00695c", "MTN Zambia": "#ffcc00", "Airtel Zambia": "#e53935"},
            )
            st.plotly_chart(fig_op_churn, use_container_width=True)
        with oc2:
            fig_op_rev = px.bar(
                op_kpis, x="operator", y="monthly_revenue", text_auto=".2s",
                title="Monthly Revenue by Operator (K)",
                color="operator",
                color_discrete_map={"Zamtel": "#00695c", "MTN Zambia": "#ffcc00", "Airtel Zambia": "#e53935"},
            )
            st.plotly_chart(fig_op_rev, use_container_width=True)
        if FOCUS_OPERATOR in op_kpis["operator"].values:
            zam = op_kpis.loc[op_kpis["operator"] == FOCUS_OPERATOR].iloc[0]
            _insight(
                f"**{FOCUS_OPERATOR}** portfolio: {int(zam['subscribers']):,} subscribers · "
                f"ARPU {_fmt_kwacha(zam['avg_arpu'])} · churn {zam.get('churn_rate_pct', 0):.1f}% · "
                f"share {zam['market_share_pct']:.1f}% — see **Step 10** for full benchmarking."
            )
            ss.eda_insights.append(
                f"{FOCUS_OPERATOR}: {int(zam['subscribers']):,} subs, ARPU {_fmt_kwacha(zam['avg_arpu'])}, "
                f"churn {zam.get('churn_rate_pct', 0):.1f}%"
            )

    if st.button("Complete EDA & Proceed →", type="primary"):
        _complete_step(2)
        st.rerun()


# ─────────────────────────────────────────────────────────────────────────────
# STEP 3 — DATA CLEANING
# ─────────────────────────────────────────────────────────────────────────────
def step_clean() -> None:
    _step_header(
        3, "Data Cleaning",
        "Remove duplicates, normalise bad tokens, impute missing values, and coerce types.",
    )
    if ss.raw_df is None:
        st.warning("Complete Step 1 first.")
        return

    df = ss.raw_df.copy()
    log: list[str] = []

    st.subheader("Cleaning Options")
    c1, c2, c3, c4 = st.columns(4)
    drop_dup = c1.checkbox("Remove duplicate rows", value=True)
    fix_null_tokens = c2.checkbox("Normalise null tokens (N/A, null, None)", value=True)
    strip_blanks = c3.checkbox("Strip blanks → null", value=True)
    fix_types = c4.checkbox("Auto-coerce numeric columns", value=True)
    impute = st.checkbox("Impute missing numeric values (median)", value=True)

    if st.button("Run Cleaning Pipeline", type="primary"):
        # 1. Strip whitespace on object columns
        if strip_blanks:
            for c in df.select_dtypes("object").columns:
                df[c] = df[c].astype(str).str.strip()
                blanks = df[c].isin(["", "nan", "NaN"])
                if blanks.any():
                    df.loc[blanks, c] = np.nan
                    log.append(f"Blank strings in `{c}` → null ({int(blanks.sum())} cells).")

        # 2. Null-token normalisation
        if fix_null_tokens:
            null_tokens = {"N/A", "n/a", "NA", "null", "NULL", "None", "none", "-"}
            for c in df.select_dtypes("object").columns:
                mask = df[c].isin(null_tokens)
                if mask.any():
                    df.loc[mask, c] = np.nan
                    log.append(f"Null tokens in `{c}` normalised ({int(mask.sum())} cells).")

        # 3. Duplicates
        if drop_dup:
            before = len(df)
            df = df.drop_duplicates().reset_index(drop=True)
            removed = before - len(df)
            if removed:
                log.append(f"Removed {removed} duplicate rows.")

        # 4. Impute numerics (first pass)
        num_cols = df.select_dtypes(include="number").columns
        if impute and len(num_cols):
            imputer = SimpleImputer(strategy="median")
            df[num_cols] = imputer.fit_transform(df[num_cols])
            log.append(f"Median-imputed {len(num_cols)} numeric columns.")

        # 5. Coerce bad numeric tokens (e.g. monthly_charges '-')
        if fix_types:
            for c in df.select_dtypes("object").columns:
                if c == ss.target_col:
                    continue
                converted = pd.to_numeric(df[c], errors="coerce")
                if converted.notnull().mean() > 0.9:
                    df[c] = converted
                    log.append(f"Auto-coerced `{c}` to numeric.")

        # 6. Re-impute numerics after coercion
        num_cols = df.select_dtypes(include="number").columns
        if impute and len(num_cols) and df[num_cols].isnull().any().any():
            imputer = SimpleImputer(strategy="median")
            df[num_cols] = imputer.fit_transform(df[num_cols])
            log.append("Re-imputed numeric columns after type coercion.")

        ss.cleaned_df = df
        ss.clean_log = log
        _ok(f"Cleaning complete — {len(df):,} rows remain.")

    if ss.cleaned_df is not None:
        st.subheader("Cleaning Log")
        for entry in ss.clean_log:
            _ok(entry)

        st.subheader("Before vs After")
        raw = ss.raw_df
        clean = ss.cleaned_df
        compare = pd.DataFrame({
            "Metric": ["Rows", "Columns", "Missing Cells", "Duplicates"],
            "Before": [
                raw.shape[0], raw.shape[1],
                int(raw.isnull().sum().sum()), int(raw.duplicated().sum()),
            ],
            "After": [
                clean.shape[0], clean.shape[1],
                int(clean.isnull().sum().sum()), int(clean.duplicated().sum()),
            ],
        })
        compare["Delta"] = compare["After"] - compare["Before"]
        st.dataframe(compare, use_container_width=True, hide_index=True)

        st.subheader("Cleaned Data Preview")
        st.dataframe(ss.cleaned_df.head(15), use_container_width=True)
        if st.button("Confirm Cleaning & Proceed →", type="primary"):
            _complete_step(3)
            st.rerun()
    else:
        st.info("Click **Run Cleaning Pipeline** to process the data.")


# ─────────────────────────────────────────────────────────────────────────────
# STEP 4 — FEATURE ENGINEERING
# ─────────────────────────────────────────────────────────────────────────────
def step_features() -> None:
    _step_header(
        4, "Feature Engineering",
        "Derived features, encoding, and scaling are applied automatically. "
        "Select which columns to include as model features.",
    )
    src = ss.cleaned_df if ss.cleaned_df is not None else ss.raw_df
    if src is None:
        st.warning("Complete Steps 1–3 first.")
        return

    df = src.copy()
    tgt = ss.target_col
    derived: list[str] = []

    df = _derive_business_kpis(df)
    if "estimated_clv" in df.columns:
        derived.append("estimated_clv = monthly_charges × tenure_months (ZMW)")
    if "rfm_recency" in df.columns:
        derived.append("rfm_recency / rfm_frequency / rfm_monetary — RFM-style scores")
    if "_target_payment_risk" in df.columns:
        rate = df["_target_payment_risk"].mean()
        derived.append(f"payment_risk target ({rate:.1%} positive)")
    if "_target_upgrade" in df.columns:
        rate = df["_target_upgrade"].mean()
        derived.append(f"upgrade_propensity target ({rate:.1%} positive)")

    st.subheader("Derived KPI Features")

    if "monthly_charges" in df.columns and "data_volume_mb" in df.columns:
        df["charge_per_mb"] = (df["monthly_charges"] / df["data_volume_mb"].replace(0, np.nan)).round(4)
        derived.append("charge_per_mb = monthly_charges / data_volume_mb")

    if "complaint_count" in df.columns and "tenure_months" in df.columns:
        df["complaint_rate"] = (
            df["complaint_count"] / df["tenure_months"].replace(0, np.nan)
        ).round(4)
        derived.append("complaint_rate = complaint_count / tenure_months")

    eng_cols = [c for c in ["call_minutes", "sms_count", "data_volume_mb"] if c in df.columns]
    if eng_cols:
        for c in eng_cols:
            std = df[c].std()
            mean = df[c].mean()
            df[f"{c}_z"] = (df[c] - mean) / (std if std > 0 else 1)
        df["engagement_score"] = df[[f"{c}_z" for c in eng_cols]].mean(axis=1).round(4)
        df = df.drop(columns=[f"{c}_z" for c in eng_cols])
        derived.append("engagement_score = mean z-score of call/SMS/data usage")

    if "tenure_months" in df.columns:
        df["tenure_bucket"] = pd.cut(
            df["tenure_months"], bins=[0, 6, 12, 24, 48, 999],
            labels=["0-6m", "6-12m", "12-24m", "24-48m", "48m+"],
        ).astype(str)
        derived.append("tenure_bucket = binned tenure_months")

    if all(c in df.columns for c in ["created_date", "closed_date"]):
        created = pd.to_datetime(df["created_date"], errors="coerce")
        closed = pd.to_datetime(df["closed_date"], errors="coerce")
        df["resolution_hours"] = ((closed - created).dt.total_seconds() / 3600).round(2)
        if "sla_hours" in df.columns:
            df["sla_status"] = np.where(
                df["resolution_hours"] <= df["sla_hours"], "Within SLA", "Breached SLA"
            )
            derived.append("sla_status + resolution_hours from ticket timestamps")

    for d in derived:
        _ok(d)

    st.subheader("Categorical Encoding")
    cat_cols = df.select_dtypes("object").columns.tolist()
    cat_cols = [c for c in cat_cols if c not in _SKIP_ENCODE_COLS]

    if cat_cols:
        for c in cat_cols:
            if c == tgt:
                continue
            n_unique = df[c].nunique()
            if n_unique > _MAX_OHE_CARDINALITY:
                _warn(
                    f"Skipped `{c}` — {n_unique} unique values "
                    f"(max {_MAX_OHE_CARDINALITY} for one-hot encoding)."
                )
                continue
            if n_unique <= 2:
                le = LabelEncoder()
                df[c] = le.fit_transform(df[c].astype(str))
                _ok(f"Label-encoded `{c}` ({n_unique} classes)")
            else:
                dummies = pd.get_dummies(df[c], prefix=c, drop_first=True)
                df = pd.concat([df.drop(columns=[c]), dummies], axis=1)
                _ok(f"One-hot-encoded `{c}` → {len(dummies.columns)} dummy columns")

    for skip in _SKIP_ENCODE_COLS:
        if skip in df.columns and skip != tgt:
            _warn(f"Excluded `{skip}` from encoding (identifier / date field).")

    st.subheader("Select Model Features")
    exclude = {"subscriber_id", "msisdn", tgt, "_target_payment_risk", "_target_upgrade"} | _SKIP_ENCODE_COLS
    candidates = [c for c in df.columns if c not in exclude and c != tgt]
    default_sel = [c for c in candidates if c in ss.feature_cols] or candidates[:20]
    selected = st.multiselect("Features to include:", candidates, default=default_sel)

    st.subheader("Feature Scaling")
    scale_method = st.selectbox("Scaling method:", ["StandardScaler", "None"])

    if st.button("Apply Features & Split Data", type="primary"):
        target = _prepare_binary_target(df[tgt].copy(), tgt)
        if target is None:
            return

        valid = target.notna()
        feat_df = df.loc[valid, selected].copy()
        target = target.loc[valid]

        X_train, X_test, y_train, y_test = _safe_train_test_split(feat_df, target)

        ss.imputer = None
        ss.scaler = None
        if X_train.isnull().any().any() or X_test.isnull().any().any():
            imputer = SimpleImputer(strategy="median")
            cols = X_train.columns.tolist()
            X_train = pd.DataFrame(imputer.fit_transform(X_train), columns=cols, index=X_train.index)
            X_test = pd.DataFrame(imputer.transform(X_test), columns=cols, index=X_test.index)
            ss.imputer = imputer
            _ok("Imputed remaining NaNs (fit on train only).")

        if scale_method == "StandardScaler":
            scaler = StandardScaler()
            cols = X_train.columns.tolist()
            X_train = pd.DataFrame(scaler.fit_transform(X_train), columns=cols, index=X_train.index)
            X_test = pd.DataFrame(scaler.transform(X_test), columns=cols, index=X_test.index)
            ss.scaler = scaler
            _ok("Applied StandardScaler (fit on train, transform test).")

        # Secondary targets aligned to the same split
        ss.y_train_payment = ss.y_test_payment = None
        ss.y_train_upgrade = ss.y_test_upgrade = None
        if "_target_payment_risk" in df.columns:
            ypay = _prepare_binary_target(df.loc[valid, "_target_payment_risk"], "_target_payment_risk")
            if ypay is not None:
                ss.y_train_payment = ypay.loc[X_train.index]
                ss.y_test_payment = ypay.loc[X_test.index]
        if "_target_upgrade" in df.columns:
            yup = _prepare_binary_target(df.loc[valid, "_target_upgrade"], "_target_upgrade")
            if yup is not None:
                ss.y_train_upgrade = yup.loc[X_train.index]
                ss.y_test_upgrade = yup.loc[X_test.index]

        ss.featured_df = df
        ss.feature_cols = selected
        ss.X_train = X_train
        ss.X_test = X_test
        ss.y_train = y_train
        ss.y_test = y_test
        ss.secondary_models = {}
        ss.secondary_results = {}
        _ok(
            f"Train: {X_train.shape[0]:,} rows · Test: {X_test.shape[0]:,} rows · "
            f"{len(selected)} features"
        )

    if ss.X_train is not None:
        st.subheader("Train / Test Split Summary")
        split_info = pd.DataFrame({
            "Set": ["Train", "Test"],
            "Rows": [len(ss.X_train), len(ss.X_test)],
            "Features": [ss.X_train.shape[1], ss.X_test.shape[1]],
            "Churn Rate": [
                f"{ss.y_train.mean():.1%}",
                f"{ss.y_test.mean():.1%}",
            ],
        })
        st.dataframe(split_info, use_container_width=True, hide_index=True)

        st.subheader("Feature Matrix Preview (Train)")
        st.dataframe(ss.X_train.head(10), use_container_width=True)

        st.subheader("Feature Statistics (Train)")
        st.dataframe(ss.X_train.describe().T.round(3), use_container_width=True, height=320)

        if st.button("Confirm Features & Proceed →", type="primary"):
            _complete_step(4)
            st.rerun()






# ─────────────────────────────────────────────────────────────────────────────
# STEP 5 — MODEL TRAINING
# ─────────────────────────────────────────────────────────────────────────────
def step_train() -> None:
    _step_header(5, "Model Training")
    st.markdown(
        "Six models are trained and compared: "
        "**Logistic Regression**, **Random Forest**, **Gradient Boosting**, "
        "**Support Vector Machine**, **K-Nearest Neighbors**, and **AdaBoost**."
    )

    if ss.X_train is None:
        st.warning("Complete Step 4 first.")
        return

    with st.expander("Hyperparameters", expanded=False):
        col1, col2, col3 = st.columns(3)
        lr_c = col1.slider("LR: regularisation C", 0.01, 10.0, 1.0, 0.01)
        rf_trees = col2.slider("RF: n_estimators", 50, 500, 200, 50)
        gb_lr = col3.slider("GB: learning_rate", 0.01, 0.5, 0.1, 0.01)
        gb_trees = col3.slider("GB: n_estimators", 50, 300, 100, 50)
        class_wt = col1.checkbox("Use class_weight='balanced'", value=True)

        col4, col5, col6 = st.columns(3)
        svm_c = col4.slider("SVM: regularisation C", 0.1, 10.0, 1.0, 0.1)
        knn_k = col5.slider("KNN: neighbours (k)", 3, 25, 5, 1)
        ada_trees = col6.slider("AdaBoost: n_estimators", 50, 300, 100, 50)

    st.markdown(
        """
| Model | Type | Strengths |
|-------|------|-----------|
| Logistic Regression | Linear | Fast baseline, interpretable coefficients |
| Random Forest | Ensemble (bagging) | Non-linear patterns, feature importance |
| Gradient Boosting | Ensemble (boosting) | Strong tabular performance |
| Support Vector Machine | Kernel classifier | Effective in high-dimensional space |
| K-Nearest Neighbors | Instance-based | Simple, no training phase |
| AdaBoost | Boosting | Combines weak learners adaptively |
"""
    )

    if st.button("Train All Models", type="primary"):
        X_train = ss.X_train.copy()
        X_test = ss.X_test.copy()
        y_train = ss.y_train.copy()
        y_test = ss.y_test.copy()

        if X_train.isnull().any().any() or X_test.isnull().any().any():
            imputer = SimpleImputer(strategy="median")
            cols = X_train.columns.tolist()
            X_train = pd.DataFrame(imputer.fit_transform(X_train), columns=cols, index=X_train.index)
            X_test = pd.DataFrame(imputer.transform(X_test), columns=cols, index=X_test.index)
            _warn("NaN values imputed before training (median, fit on train).")

        cw = "balanced" if class_wt else None
        model_defs = {
            "Logistic Regression": LogisticRegression(C=lr_c, max_iter=1000, class_weight=cw),
            "Random Forest": RandomForestClassifier(
                n_estimators=rf_trees, random_state=42, class_weight=cw, n_jobs=-1
            ),
            "Gradient Boosting": GradientBoostingClassifier(
                n_estimators=gb_trees, learning_rate=gb_lr, random_state=42
            ),
            "Support Vector Machine": SVC(C=svm_c, probability=True, class_weight=cw, random_state=42),
            "K-Nearest Neighbors": KNeighborsClassifier(n_neighbors=knn_k),
            "AdaBoost": AdaBoostClassifier(n_estimators=ada_trees, random_state=42),
        }

        results = {}
        trained = {}
        with st.spinner("Training models…"):
            for name, model in model_defs.items():
                st.write(f"Training {name}…")
                model.fit(X_train, y_train)
                y_pred = model.predict(X_test)
                y_prob = (
                    model.predict_proba(X_test)[:, 1]
                    if hasattr(model, "predict_proba") else y_pred.astype(float)
                )
                results[name] = {
                    "accuracy": accuracy_score(y_test, y_pred),
                    "precision": precision_score(y_test, y_pred, zero_division=0),
                    "recall": recall_score(y_test, y_pred, zero_division=0),
                    "f1": f1_score(y_test, y_pred, zero_division=0),
                    "roc_auc": roc_auc_score(y_test, y_prob),
                    "y_pred": y_pred,
                    "y_prob": y_prob,
                }
                trained[name] = model

        ss.models = trained
        ss.model_results = results

        if "Random Forest" in trained:
            rf = trained["Random Forest"]
            imp = pd.Series(rf.feature_importances_, index=X_train.columns).sort_values(ascending=False)
            ss.feature_importance = imp

        _ok(f"{len(trained)} churn models trained successfully!")

        # ── Secondary models: payment risk & upgrade propensity ───────────────
        st.subheader("Secondary Predictions")
        sec_models = {}
        sec_results = {}
        if ss.y_train_payment is not None and ss.y_test_payment is not None:
            mdl, met = _train_secondary_rf(
                X_train, X_test, ss.y_train_payment, ss.y_test_payment, "payment_risk"
            )
            if mdl and met:
                sec_models["Payment Risk"] = mdl
                sec_results["Payment Risk"] = met
                _ok(f"Payment Risk model — ROC-AUC {met['roc_auc']:.3f} (positive rate {met['positive_rate']:.1%})")
        if ss.y_train_upgrade is not None and ss.y_test_upgrade is not None:
            mdl, met = _train_secondary_rf(
                X_train, X_test, ss.y_train_upgrade, ss.y_test_upgrade, "upgrade"
            )
            if mdl and met:
                sec_models["Upgrade Propensity"] = mdl
                sec_results["Upgrade Propensity"] = met
                _ok(f"Upgrade Propensity model — ROC-AUC {met['roc_auc']:.3f} (positive rate {met['positive_rate']:.1%})")
        ss.secondary_models = sec_models
        ss.secondary_results = sec_results

    if ss.model_results:
        st.subheader("Training Summary")
        summary = pd.DataFrame({
            name: {k: v for k, v in res.items() if k not in ("y_pred", "y_prob")}
            for name, res in ss.model_results.items()
        }).T.round(4)
        st.dataframe(summary, use_container_width=True)

        best = max(ss.model_results, key=lambda n: ss.model_results[n]["roc_auc"])
        _insight(f"Best model by ROC-AUC: **{best}** ({ss.model_results[best]['roc_auc']:.3f})")

        if ss.feature_importance is not None:
            st.subheader("Top Feature Importances (Random Forest)")
            fi = ss.feature_importance.head(15).reset_index()
            fi.columns = ["Feature", "Importance"]
            fig = px.bar(fi, x="Importance", y="Feature", orientation="h", title="Top 15 Features")
            st.plotly_chart(fig, use_container_width=True)

        if ss.secondary_results:
            st.subheader("Secondary Model Scorecard")
            sec_df = pd.DataFrame(ss.secondary_results).T.round(4)
            st.dataframe(sec_df, use_container_width=True)
            _insight(
                "Payment Risk → proactive billing outreach · "
                "Upgrade Propensity → upsell 5G/premium plans on stable subscribers."
            )

        if st.button("Confirm Training & Proceed →", type="primary"):
            _complete_step(5)
            st.rerun()


# ─────────────────────────────────────────────────────────────────────────────
# STEP 6 — MODEL EVALUATION
# ─────────────────────────────────────────────────────────────────────────────
def step_evaluate() -> None:
    _step_header(6, "Model Evaluation", "Compare confusion matrices, ROC curves, and classification reports.")

    if not ss.model_results:
        st.warning("Complete Step 5 first.")
        return

    y_test = ss.y_test
    model_names = list(ss.model_results.keys())
    pick = st.selectbox("Select model for detailed view:", model_names)
    res = ss.model_results[pick]

    st.markdown(f"### {pick} — Performance Card")
    card_cols = st.columns(6)
    card_cols[0].metric("Accuracy", f"{res['accuracy']:.1%}")
    card_cols[1].metric("Precision", f"{res['precision']:.1%}")
    card_cols[2].metric("Recall", f"{res['recall']:.1%}")
    card_cols[3].metric("F1 Score", f"{res['f1']:.1%}")
    card_cols[4].metric("ROC-AUC", f"{res['roc_auc']:.3f}")
    rank = sorted(ss.model_results, key=lambda n: ss.model_results[n]["roc_auc"], reverse=True).index(pick) + 1
    card_cols[5].metric("AUC Rank", f"#{rank} / {len(model_names)}")

    st.subheader("Side-by-Side Model Ranking")
    rank_df = pd.DataFrame({
        name: {
            "Accuracy": res["accuracy"],
            "Precision": res["precision"],
            "Recall": res["recall"],
            "F1": res["f1"],
            "ROC-AUC": res["roc_auc"],
        }
        for name, res in ss.model_results.items()
    }).T.sort_values("ROC-AUC", ascending=False).round(4)
    rank_df["Rank"] = range(1, len(rank_df) + 1)
    st.dataframe(rank_df, use_container_width=True)

    fig_rank = px.bar(
        rank_df.reset_index().rename(columns={"index": "Model"}),
        x="ROC-AUC", y="Model", orientation="h", color="ROC-AUC",
        color_continuous_scale="Blues", title="Models Ranked by ROC-AUC",
    )
    st.plotly_chart(fig_rank, use_container_width=True)

    c1, c2, c3, c4, c5 = st.columns(5)
    c1.metric("Accuracy", f"{res['accuracy']:.3f}")
    c2.metric("Precision", f"{res['precision']:.3f}")
    c3.metric("Recall", f"{res['recall']:.3f}")
    c4.metric("F1", f"{res['f1']:.3f}")
    c5.metric("ROC-AUC", f"{res['roc_auc']:.3f}")

    col_cm, col_roc = st.columns(2)
    with col_cm:
        st.subheader("Confusion Matrix")
        cm = confusion_matrix(y_test, res["y_pred"])
        fig_cm = px.imshow(
            cm, text_auto=True, x=["Pred 0", "Pred 1"], y=["Actual 0", "Actual 1"],
            color_continuous_scale="Blues", title=f"{pick} — Confusion Matrix",
        )
        st.plotly_chart(fig_cm, use_container_width=True)

    with col_roc:
        st.subheader("ROC Curve — All Models")
        fig_roc = go.Figure()
        for name, r in ss.model_results.items():
            fpr, tpr, _ = roc_curve(y_test, r["y_prob"])
            fig_roc.add_trace(go.Scatter(
                x=fpr, y=tpr, mode="lines", name=f"{name} (AUC={r['roc_auc']:.3f})"
            ))
        fig_roc.add_trace(go.Scatter(x=[0, 1], y=[0, 1], mode="lines",
                                     line=dict(dash="dash", color="gray"), name="Random"))
        fig_roc.update_layout(title="ROC Curves", xaxis_title="FPR", yaxis_title="TPR")
        st.plotly_chart(fig_roc, use_container_width=True)

    st.subheader("Classification Report")
    report = classification_report(y_test, res["y_pred"], output_dict=True, zero_division=0)
    st.dataframe(pd.DataFrame(report).T.round(3), use_container_width=True)

    st.subheader("Churn Probability Distribution")
    prob_df = pd.DataFrame({"churn_probability": res["y_prob"]})
    fig_hist = px.histogram(prob_df, x="churn_probability", nbins=30, title="P(churn) on Test Set")
    st.plotly_chart(fig_hist, use_container_width=True)

    st.subheader("Model Comparison Heatmap")
    metrics = ["accuracy", "precision", "recall", "f1", "roc_auc"]
    hm = pd.DataFrame({
        name: [res[m] for m in metrics]
        for name, res in ss.model_results.items()
    }, index=metrics)
    fig_hm = px.imshow(
        hm, text_auto=".3f", aspect="auto",
        color_continuous_scale="Viridis", title="Metrics Heatmap — All Models",
    )
    st.plotly_chart(fig_hm, use_container_width=True)

    if st.button("Complete Evaluation & Proceed →", type="primary"):
        _complete_step(6)
        st.rerun()


# ─────────────────────────────────────────────────────────────────────────────
# STEP 7 — BIAS DETECTION
# ─────────────────────────────────────────────────────────────────────────────
def step_bias() -> None:
    _step_header(
        7, "Bias Detection",
        "**Learning curves** show training score vs. cross-validation score as the training size grows.  \n"
        "- **Overfit**: high train score, low val score (gap stays large)  \n"
        "- **Underfit**: both scores are low and close together  \n"
        "- **Good fit**: both scores converge at a high value",
    )

    if not ss.models or ss.X_train is None:
        st.warning("Complete Step 5 first.")
        return

    model_names = list(ss.models.keys())
    pick = st.selectbox("Model for learning curve:", model_names)
    model = ss.models[pick]

    if st.button("Compute Learning Curve", type="primary"):
        with st.spinner("Computing learning curve (may take a moment)…"):
            train_sizes, train_scores, val_scores = learning_curve(
                model, ss.X_train, ss.y_train, cv=3,
                train_sizes=np.linspace(0.2, 1.0, 5), scoring="roc_auc", n_jobs=-1,
            )
        train_mean = train_scores.mean(axis=1)
        val_mean = val_scores.mean(axis=1)

        fig = go.Figure()
        fig.add_trace(go.Scatter(x=train_sizes, y=train_mean, mode="lines+markers", name="Train AUC"))
        fig.add_trace(go.Scatter(x=train_sizes, y=val_mean, mode="lines+markers", name="CV AUC"))
        fig.update_layout(title=f"Learning Curve — {pick}", xaxis_title="Training Size", yaxis_title="ROC-AUC")
        st.plotly_chart(fig, use_container_width=True)

        gap = train_mean[-1] - val_mean[-1]
        if gap > 0.08:
            diagnosis = "Overfitting"
        elif val_mean[-1] < 0.65:
            diagnosis = "Underfitting"
        else:
            diagnosis = "Good Fit"
        _insight(f"Diagnosis for **{pick}**: **{diagnosis}** (train–CV gap = {gap:.3f})")

    st.subheader("Train vs Test AUC Gap")
    rows = []
    for name, model in ss.models.items():
        if not hasattr(model, "predict_proba"):
            continue
        train_auc = roc_auc_score(ss.y_train, model.predict_proba(ss.X_train)[:, 1])
        test_auc = ss.model_results[name]["roc_auc"]
        gap = train_auc - test_auc
        if gap > 0.08:
            fit = "Overfitting"
        elif test_auc < 0.65:
            fit = "Underfitting"
        else:
            fit = "Good Fit"
        rows.append({"Model": name, "Train AUC": train_auc, "Test AUC": test_auc, "Gap": gap, "Fit": fit})

    if rows:
        gap_df = pd.DataFrame(rows).round(4)
        st.dataframe(gap_df, use_container_width=True, hide_index=True)
        fig_gap = px.bar(gap_df, x="Model", y="Gap", color="Fit", title="Train − Test AUC Gap")
        st.plotly_chart(fig_gap, use_container_width=True)

    if st.button("Complete Bias Check & Proceed →", type="primary"):
        _complete_step(7)
        st.rerun()


# ─────────────────────────────────────────────────────────────────────────────
# STEP 8 — RECOMMENDATIONS
# ─────────────────────────────────────────────────────────────────────────────
SEGMENT_CAMPAIGNS = {
    "High-Risk / Service Issue": {
        "product": "Priority Support Plan",
        "campaign": "Service Recovery — Personal apology + 3-month fee waiver",
        "urgency": "CRITICAL",
        "color": "reco-red",
    },
    "High-Risk / Low Usage": {
        "product": "Starter Data Bundle 2 GB @ K45",
        "campaign": "Win-Back — Free 1GB data for 2 months",
        "urgency": "HIGH",
        "color": "reco-amber",
    },
    "High-Risk / High Value": {
        "product": "Loyalty Rewards Programme",
        "campaign": "Retention — 20% discount + dedicated account manager",
        "urgency": "HIGH",
        "color": "reco-amber",
    },
    "At-Risk": {
        "product": "Flexi Bundle (voice + data + SMS)",
        "campaign": "Proactive Engagement — Usage milestone reward",
        "urgency": "MEDIUM",
        "color": "reco-amber",
    },
    "New Subscriber": {
        "product": "Welcome Pack + 30-day trial upgrade",
        "campaign": "Onboarding Journey — Guided tutorials + first-month bonus",
        "urgency": "MEDIUM",
        "color": "reco-blue",
    },
    "Loyal / High Value": {
        "product": "5G Early Access / Premium Plan Upgrade",
        "campaign": "Upsell — Early 5G access + referral bonus",
        "urgency": "LOW",
        "color": "reco-green",
    },
    "Stable / Standard": {
        "product": "Data Add-On or SMS Bundle",
        "campaign": "Cross-Sell — Personalised data top-up offers",
        "urgency": "LOW",
        "color": "reco-green",
    },
}


def step_recommendations() -> None:
    _step_header(
        8, "Recommendations",
        "Score every subscriber for churn probability and assign one of seven actionable segments.",
    )

    if not ss.model_results or ss.X_test is None:
        st.warning("Complete Steps 5–6 first.")
        return

    best = max(ss.model_results, key=lambda n: ss.model_results[n]["roc_auc"])
    model = ss.models[best]

    def _segment(row) -> str:
        p = row["churn_probability"]
        complaints = row.get("complaint_count", 0) or 0
        data_mb = row.get("data_volume_mb", 0) or 0
        charges = row.get("monthly_charges", 0) or 0
        tenure = row.get("tenure_months", 12) or 12

        if p >= 0.65 and complaints >= 3:
            return "High-Risk / Service Issue"
        if p >= 0.65 and data_mb < 200:
            return "High-Risk / Low Usage"
        if p >= 0.65 and charges > HIGH_VALUE_KWACHA:
            return "High-Risk / High Value"
        if p >= 0.4:
            return "At-Risk"
        if tenure < 6:
            return "New Subscriber"
        if data_mb > 3000 or charges > LOYAL_VALUE_KWACHA:
            return "Loyal / High Value"
        return "Stable / Standard"

    if st.button("Generate Recommendations", type="primary"):
        # Score full subscriber base when featured matrix is available
        if ss.featured_df is not None and ss.feature_cols:
            score_df = ss.featured_df[ss.feature_cols].copy()
            score_df = _apply_feature_pipeline(score_df)
            probs = model.predict_proba(score_df)[:, 1]
            reco = ss.featured_df.loc[score_df.index].copy()
            reco["churn_probability"] = probs.round(4)
        else:
            score_df = ss.X_test.copy()
            probs = model.predict_proba(score_df)[:, 1]
            reco = score_df.copy()
            reco["churn_probability"] = probs.round(4)

        src = ss.cleaned_df if ss.cleaned_df is not None else ss.raw_df
        meta_cols = [
            c for c in [
                "complaint_count", "data_volume_mb", "monthly_charges", "tenure_months",
                "estimated_clv", "payment_failures", "days_overdue", "bill_shock_flag",
            ] if src is not None and c in src.columns
        ]
        for c in meta_cols:
            if c in reco.columns:
                continue
            if c in score_df.columns:
                reco[c] = score_df[c].values
            elif src is not None:
                aligned = src.reindex(reco.index)
                reco[c] = aligned[c].values

        charges = pd.to_numeric(reco.get("monthly_charges", 0), errors="coerce").fillna(0)
        reco["revenue_at_risk"] = (charges * reco["churn_probability"]).round(2)
        if "estimated_clv" not in reco.columns and "monthly_charges" in reco.columns and "tenure_months" in reco.columns:
            reco["estimated_clv"] = (charges * pd.to_numeric(reco["tenure_months"], errors="coerce").fillna(0)).round(2)

        # Secondary model scores
        if ss.secondary_models and ss.feature_cols:
            base_X = (
                ss.featured_df.loc[reco.index, ss.feature_cols].copy()
                if ss.featured_df is not None else ss.X_test.copy()
            )
            sec_X = _apply_feature_pipeline(base_X)
            if "Payment Risk" in ss.secondary_models:
                reco["payment_risk_probability"] = ss.secondary_models["Payment Risk"].predict_proba(sec_X)[:, 1].round(4)
            if "Upgrade Propensity" in ss.secondary_models:
                reco["upgrade_probability"] = ss.secondary_models["Upgrade Propensity"].predict_proba(sec_X)[:, 1].round(4)

        reco["segment"] = reco.apply(_segment, axis=1)
        reco["action_segment"] = reco.apply(_action_segment, axis=1)
        for field in ("product", "campaign", "urgency", "color"):
            reco[field] = reco["segment"].map(lambda s, f=field: SEGMENT_CAMPAIGNS[s][f])

        ss.recommendations_df = reco
        total_rar = float(reco["revenue_at_risk"].sum())
        high_risk = int((reco["churn_probability"] >= 0.5).sum())
        ss.analytics_kpis.update({
            "total_revenue_at_risk": total_rar,
            "high_risk_subscribers": high_risk,
            "avg_clv": float(reco["estimated_clv"].mean()) if "estimated_clv" in reco.columns else None,
        })
        _ok(
            f"Scored {len(reco):,} subscribers — {high_risk:,} high-risk · "
            f"total revenue at risk **{_fmt_kwacha(total_rar)}**"
        )

    if ss.recommendations_df is not None:
        reco = ss.recommendations_df

        st.subheader("Revenue at Risk Dashboard")
        rar_cols = st.columns(4)
        total_rar = reco["revenue_at_risk"].sum()
        rar_cols[0].metric("Total Revenue at Risk", _fmt_kwacha(total_rar))
        rar_cols[1].metric("Avg RAR / Subscriber", _fmt_kwacha(reco["revenue_at_risk"].mean()))
        save_n = int((reco["action_segment"] == "Save — High Value at Risk").sum()) if "action_segment" in reco.columns else 0
        rar_cols[2].metric("Save List (high-value)", f"{save_n:,}")
        if "estimated_clv" in reco.columns:
            rar_cols[3].metric("Avg Estimated CLV", _fmt_kwacha(reco["estimated_clv"].mean()))

        if "segment" in reco.columns:
            seg_rar = (
                reco.groupby("segment")["revenue_at_risk"].sum()
                .sort_values(ascending=False).reset_index()
            )
            seg_rar.columns = ["Segment", "Revenue at Risk"]
            fig_rar = px.bar(
                seg_rar, x="Segment", y="Revenue at Risk",
                color="Revenue at Risk", color_continuous_scale="Reds",
                title="Revenue at Risk by Churn Segment",
            )
            st.plotly_chart(fig_rar, use_container_width=True)

        if "action_segment" in reco.columns:
            st.subheader("Action Segments — Save / Grow / Fix / Monitor")
            act = reco["action_segment"].value_counts().reset_index()
            act.columns = ["Action", "Subscribers"]
            fig_act = px.pie(act, names="Action", values="Subscribers", title="Strategic Action Mix")
            st.plotly_chart(fig_act, use_container_width=True)
            for action in act["Action"]:
                sub = reco[reco["action_segment"] == action]
                hint = {
                    "Save — High Value at Risk": "Immediate retention offer + dedicated agent",
                    "Fix — Billing & Payment Risk": "Payment plan, bill explain, fee waiver review",
                    "Grow — Upsell Ready": "5G/premium plan cross-sell campaign",
                    "Watch — Elevated Churn Risk": "Proactive engagement + usage nudge",
                    "Monitor — Stable": "Standard cross-sell only",
                }.get(action, "Review segment rules")
                _insight(f"**{action}** ({len(sub):,} subs): {hint}")

        st.subheader("Customer Segmentation")
        seg_counts = reco["segment"].value_counts().reset_index()
        seg_counts.columns = ["Segment", "Count"]

        c1, c2 = st.columns(2)
        with c1:
            fig_pie = px.pie(seg_counts, names="Segment", values="Count", title="Segment Distribution")
            st.plotly_chart(fig_pie, use_container_width=True)
        with c2:
            seg_avg = reco.groupby("segment")["churn_probability"].mean().reset_index()
            seg_avg.columns = ["segment", "Avg P(churn)"]
            fig_bar = px.bar(
                seg_avg.sort_values("Avg P(churn)", ascending=True),
                x="Avg P(churn)", y="segment", orientation="h",
                color="Avg P(churn)", color_continuous_scale="Reds",
                title="Average Churn Probability by Segment",
            )
            st.plotly_chart(fig_bar, use_container_width=True)

        st.subheader("Scored Subscribers Preview")
        preview_cols = [
            c for c in [
                "churn_probability", "revenue_at_risk", "estimated_clv",
                "payment_risk_probability", "upgrade_probability",
                "segment", "action_segment", "complaint_count",
                "data_volume_mb", "monthly_charges", "tenure_months",
                "product", "campaign", "urgency",
            ] if c in reco.columns
        ]
        st.dataframe(
            reco[preview_cols].sort_values("churn_probability", ascending=False).head(25),
            use_container_width=True,
            hide_index=True,
        )

        st.subheader("Campaign Recommendations by Segment")
        for seg in seg_counts["Segment"]:
            sub = reco[reco["segment"] == seg]
            camp = SEGMENT_CAMPAIGNS[seg]
            st.markdown(
                f'<div class="reco-card {camp["color"]}"><strong>{seg}</strong> &nbsp;·&nbsp; '
                f'{len(sub):,} subscribers  <br/>📦 <em>Product:</em> {camp["product"]}'
                f'  <br/>📢 <em>Campaign:</em> {camp["campaign"]}'
                f'  <br/>🔴 <em>Urgency:</em> {camp["urgency"]}</div>',
                unsafe_allow_html=True,
            )

        st.download_button(
            "Download Recommendation CSV",
            data=reco.to_csv(index=False).encode(),
            file_name="churn_recommendations.csv",
            mime="text/csv",
        )

        if st.button("Proceed to Data Story →", type="primary"):
            _complete_step(8)
            st.rerun()






# ─────────────────────────────────────────────────────────────────────────────
# POWERPOINT HELPERS
# ─────────────────────────────────────────────────────────────────────────────
def _ppt_bullet_slide(prs, title: str, bullets: list[str]) -> None:
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)


def _ppt_table_slide(prs, title: str, headers: list[str], rows: list[list]) -> None:
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    nrows = len(rows) + 1
    ncols = len(headers)
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.4 * nrows)
    table = slide.shapes.add_table(nrows, ncols, left, top, width, height).table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = str(h)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            table.cell(i, j).text = str(val)


def _generate_pptx() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    df = ss.raw_df
    clean_df = ss.cleaned_df
    tgt = ss.target_col
    best = (
        max(ss.model_results, key=lambda n: ss.model_results[n]["roc_auc"])
        if ss.model_results else None
    )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Telecom Subscriber Analytics"
    sub = slide.placeholders[1].text_frame
    sub.text = (
        f"Churn Prediction — Executive Data Story\n"
        f"{ss.dataset_name}  ·  {datetime.now().strftime('%B %Y')}"
    )

    churn_rate = None
    if df is not None and tgt in df.columns:
        churn_rate = pd.to_numeric(df[tgt], errors="coerce").mean()

    bullets = [
        "End-to-end analytics pipeline covering 9 stages:",
        f"Dataset: {ss.dataset_name}",
    ]
    if df is not None:
        bullets.append(f"{df.shape[0]:,} records, {df.shape[1]} features")
    if churn_rate is not None:
        bullets.append(f"Overall churn rate: {churn_rate:.1%}")
    if clean_df is not None:
        bullets.append(f"Clean records after preprocessing: {clean_df.shape[0]:,}")
    if ss.feature_cols:
        bullets.append(f"Model features engineered: {len(ss.feature_cols)}")
    if best:
        bullets.append(
            f"Best model: {best} (ROC-AUC = {ss.model_results[best]['roc_auc']:.3f})"
        )
    if ss.recommendations_df is not None:
        hi = int((ss.recommendations_df["churn_probability"] >= 0.5).sum())
        bullets.append(f"High-risk subscribers (P ≥ 50%): {hi:,}")
        rar = ss.analytics_kpis.get("total_revenue_at_risk")
        if rar:
            bullets.append(f"Total revenue at risk: {_fmt_kwacha(rar)}")
    _ppt_bullet_slide(prs, "Executive Summary", bullets)

    if df is not None:
        _ppt_bullet_slide(prs, "Data Quality & Preparation", [
            f"Raw records: {df.shape[0]:,}",
            f"Null rate: {df.isnull().mean().mean():.1%}",
            f"Duplicate rows: {int(df.duplicated().sum()):,}",
            f"Target column: {tgt}",
        ])

    if churn_rate is not None:
        churned = int(pd.to_numeric(df[tgt], errors="coerce").sum())
        retained = int(len(df) - churned)
        status = (
            "CRITICAL — above 25% threshold" if churn_rate > 0.25
            else "ELEVATED — proactive retention recommended" if churn_rate > 0.15
            else "Within acceptable range"
        )
        _ppt_bullet_slide(prs, "Churn Profile", [
            f"Churn rate: {churn_rate:.1%}",
            f"Churned subscribers: {churned:,}",
            f"Retained subscribers: {retained:,}",
            f"Status: {status}",
        ])

    if ss.model_results:
        headers = ["Model", "Accuracy", "Precision", "Recall", "F1", "ROC-AUC", "Fit"]
        rows = []
        for name, res in ss.model_results.items():
            train_auc = None
            if name in ss.models and hasattr(ss.models[name], "predict_proba"):
                train_auc = roc_auc_score(ss.y_train, ss.models[name].predict_proba(ss.X_train)[:, 1])
            gap = (train_auc - res["roc_auc"]) if train_auc else 0
            fit = "Overfitting" if gap > 0.08 else ("Underfitting" if res["roc_auc"] < 0.65 else "Good Fit")
            rows.append([
                name,
                f"{res['accuracy']:.3f}",
                f"{res['precision']:.3f}",
                f"{res['recall']:.3f}",
                f"{res['f1']:.3f}",
                f"{res['roc_auc']:.3f}",
                fit,
            ])
        _ppt_table_slide(prs, "Model Benchmark — Six Classifiers", headers, rows)

    if best:
        res = ss.model_results[best]
        _ppt_bullet_slide(prs, "Recommended Model", [
            f"Recommended model: {best}",
            f"Test ROC-AUC: {res['roc_auc']:.3f}",
            f"Accuracy: {res['accuracy']:.3f}  |  F1: {res['f1']:.3f}",
            f"Precision: {res['precision']:.3f}  |  Recall: {res['recall']:.3f}",
            "Deploy this model for CRM churn scoring and campaign prioritisation.",
        ])

    if ss.feature_importance is not None:
        top = ss.feature_importance.head(8)
        bullets = [f"{feat}: {imp:.4f}" for feat, imp in top.items()]
        _ppt_bullet_slide(prs, "Top Churn Drivers", bullets)

    if ss.recommendations_df is not None:
        rar = ss.analytics_kpis.get("total_revenue_at_risk", ss.recommendations_df["revenue_at_risk"].sum())
        avg_clv = ss.recommendations_df["estimated_clv"].mean() if "estimated_clv" in ss.recommendations_df.columns else None
        rar_bullets = [
            f"Total revenue at risk: {_fmt_kwacha(rar)}",
            f"Average RAR per subscriber: {_fmt_kwacha(ss.recommendations_df['revenue_at_risk'].mean())}",
        ]
        if avg_clv is not None:
            rar_bullets.append(f"Average estimated CLV: {_fmt_kwacha(avg_clv)}")
        if "action_segment" in ss.recommendations_df.columns:
            for act, cnt in ss.recommendations_df["action_segment"].value_counts().items():
                rar_bullets.append(f"{act}: {cnt:,} subscribers")
        _ppt_bullet_slide(prs, "Revenue at Risk & Action Segments", rar_bullets)

        seg = ss.recommendations_df.groupby("segment").agg(
            count=("churn_probability", "count"),
            mean=("churn_probability", "mean"),
        ).sort_values("mean", ascending=False)
        bullets = [
            f"{idx}: {int(row['count']):,} subs, avg P(churn) = {row['mean']:.2f}"
            for idx, row in seg.iterrows()
        ]
        _ppt_bullet_slide(prs, "Campaign Recommendations", bullets)

    _ppt_bullet_slide(prs, "Recommended Action Plan", [
        "Immediate (0–7 days): Contact high-risk / service-issue subscribers; escalate complaints.",
        "Short-term (1–4 weeks): Launch win-back data bundle campaign; enable CRM churn alerts.",
        "Medium-term (1–3 months): Cross-sell flexi bundles; automate onboarding journeys.",
        "Long-term (3–12 months): Premium 5G upgrades for loyal segments; real-time ML scoring.",
    ])

    if ss.model_results:
        _ppt_table_slide(
            prs,
            "Quick Model Comparison",
            ["Model", "ROC-AUC", "F1", "Accuracy"],
            [
                [
                    name,
                    f"{res['roc_auc']:.3f}",
                    f"{res['f1']:.3f}",
                    f"{res['accuracy']:.3f}",
                ]
                for name, res in sorted(
                    ss.model_results.items(),
                    key=lambda x: x[1]["roc_auc"],
                    reverse=True,
                )
            ],
        )

    if ss.recommendations_df is not None:
        urgent = ss.recommendations_df[
            ss.recommendations_df["urgency"].isin(["CRITICAL", "HIGH"])
        ]
        _ppt_bullet_slide(prs, "Priority Outreach Queue", [
            f"Critical + high urgency subscribers: {len(urgent):,}",
            f"Average P(churn) in priority queue: {urgent['churn_probability'].mean():.2f}",
            "Focus: High-Risk / Service Issue and High-Risk / High Value segments first.",
            "CRM action: assign dedicated retention agent within 48 hours.",
        ])

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Thank You"
    slide.placeholders[1].text_frame.text = (
        "Questions & Next Steps\n\n"
        "Integrate model scores into CRM  ·  Monitor monthly  ·  Retrain quarterly"
    )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _generate_technical_pptx() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    df = ss.raw_df
    clean_df = ss.cleaned_df
    tgt = ss.target_col
    best = (
        max(ss.model_results, key=lambda n: ss.model_results[n]["roc_auc"])
        if ss.model_results else None
    )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Technical Deep Dive"
    sub = slide.placeholders[1].text_frame
    sub.text = (
        f"Telco Churn ML Pipeline — Technical Documentation\n"
        f"{ss.dataset_name}  ·  {datetime.now().strftime('%B %Y')}"
    )

    _ppt_bullet_slide(prs, "Problem Definition", [
        "Business objective: predict subscriber churn and prioritise retention actions.",
        "ML task: binary classification (churn = 1, retained = 0).",
        "Constraints: interpretable features, CRM-ready scores, segment-level actions.",
        f"Dataset: {ss.dataset_name} ({df.shape[0]:,} rows)" if df is not None else "Dataset: N/A",
    ])

    _ppt_bullet_slide(prs, "Pipeline Architecture (9 Steps)", [
        "Step 1 — Data Ingestion: CSV upload / demo dataset (44 Zambia telco fields incl. operator).",
        "Step 2 — EDA: distributions, correlations, class balance, missing-value audit.",
        "Step 3 — Cleaning: dedup, null imputation, type coercion, token normalisation.",
        "Step 4 — Feature Engineering: derived KPIs, encoding, train/test split, scaling.",
        "Step 5 — Training: six classifiers with hyperparameter tuning.",
        "Step 6 — Evaluation: confusion matrix, ROC-AUC, precision/recall/F1.",
        "Step 7 — Bias Detection: learning curves, train vs test AUC gap.",
        "Step 8 — Recommendations: probability scoring + rule-based segmentation.",
        "Step 9 — Data Story: executive + technical exports.",
    ])

    if df is not None:
        _ppt_bullet_slide(prs, "Data Schema", [
            f"Columns: {df.shape[1]} (identity, subscription, usage, network, CX, billing, loyalty)",
            f"Target: {tgt}",
            "Identity keys excluded from modelling: subscriber_id, msisdn, id_number, account_number.",
            "Date fields excluded from one-hot encoding: activation_date, created_date, closed_date.",
            f"High-cardinality skip threshold: {_MAX_OHE_CARDINALITY} unique values.",
        ])

    if ss.clean_log:
        _ppt_bullet_slide(prs, "Cleaning Transformations", ss.clean_log[:12])

    _ppt_bullet_slide(prs, "Feature Engineering Details", [
        "Derived KPI — charge_per_mb = monthly_charges / data_volume_mb",
        "Derived KPI — complaint_rate = complaint_count / tenure_months",
        "Derived KPI — engagement_score = mean z-score(call_minutes, sms_count, data_volume_mb)",
        "Derived KPI — tenure_bucket = pd.cut(tenure_months, [0,6,12,24,48,999])",
        "Categorical encoding: label-encode (≤2 classes) or one-hot (>2, ≤30 cardinality)",
        "Scaling: StandardScaler (optional) fit on train, applied to test.",
        "Target validation: binary 0/1 coercion with stratified split when possible.",
    ])

    algo_notes = [
        ("Logistic Regression", "Linear baseline; fast, interpretable coefficients."),
        ("Random Forest", "Bagged trees; handles non-linearity; feature importance."),
        ("Gradient Boosting", "Sequential boosted trees; strong tabular performance."),
        ("Support Vector Machine", "Kernel margin classifier; probability via Platt scaling."),
        ("K-Nearest Neighbors", "Instance-based; local pattern matching."),
        ("AdaBoost", "Adaptive boosting of weak learners."),
    ]
    _ppt_table_slide(
        prs, "Model Zoo — Six Classifiers",
        ["Model", "Algorithm Notes"],
        [[n, d] for n, d in algo_notes],
    )

    if ss.model_results:
        headers = ["Model", "Train AUC", "Test AUC", "Gap", "Diagnosis"]
        rows = []
        for name, res in ss.model_results.items():
            train_auc = ""
            gap = ""
            diag = ""
            if name in ss.models and hasattr(ss.models[name], "predict_proba"):
                ta = roc_auc_score(ss.y_train, ss.models[name].predict_proba(ss.X_train)[:, 1])
                train_auc = f"{ta:.3f}"
                gap = f"{ta - res['roc_auc']:.3f}"
                diag = "Overfitting" if ta - res["roc_auc"] > 0.08 else (
                    "Underfitting" if res["roc_auc"] < 0.65 else "Good Fit"
                )
            rows.append([name, train_auc, f"{res['roc_auc']:.3f}", gap, diag])
        _ppt_table_slide(prs, "Train vs Test AUC Benchmark", headers, rows)

    _ppt_bullet_slide(prs, "Metrics Glossary", [
        "Accuracy — overall correct predictions.",
        "Precision — of predicted churners, how many actually churned.",
        "Recall — of actual churners, how many were caught.",
        "F1 — harmonic mean of precision and recall.",
        "ROC-AUC — rank quality of probability scores across thresholds.",
    ])

    _ppt_bullet_slide(prs, "Bias & Learning Curves", [
        "Learning curves plot train vs CV score across training sizes.",
        "Large train–CV gap → overfitting; both low → underfitting.",
        "Monitor monthly PSI on features, ROC-AUC drift, prediction volume by segment.",
        "A/B test retention campaigns by segment before full rollout.",
    ])

    if ss.secondary_results:
        headers = ["Model", "ROC-AUC", "F1", "Positive Rate"]
        rows = [
            [name, f"{m['roc_auc']:.3f}", f"{m['f1']:.3f}", f"{m['positive_rate']:.1%}"]
            for name, m in ss.secondary_results.items()
        ]
        _ppt_table_slide(prs, "Secondary Models — Payment Risk & Upgrade Propensity", headers, rows)

    _ppt_bullet_slide(prs, "Revenue at Risk & CLV Analytics", [
        "estimated_clv = monthly_charges × tenure_months",
        "revenue_at_risk = monthly_charges × P(churn) per subscriber",
        "Action segments: Save (high RAR) · Fix (payment risk) · Grow (upsell) · Monitor",
        *(
            [f"Total RAR in latest run: {_fmt_kwacha(ss.analytics_kpis['total_revenue_at_risk'])}"]
            if ss.analytics_kpis.get("total_revenue_at_risk") else
            ["Run Step 8 to compute portfolio RAR"]
        ),
    ])

    if ss.feature_importance is not None:
        headers = ["Feature", "Importance"]
        rows = [[f, f"{v:.4f}"] for f, v in ss.feature_importance.head(12).items()]
        _ppt_table_slide(prs, "Feature Importance (Random Forest — Gini Decrease)", headers, rows)

    _ppt_bullet_slide(prs, "Scoring & Segmentation Rules", [
        "Rule engine maps probability + usage KPIs to 7 segments:",
        "  P ≥ 0.65 & complaints ≥ 3  →  High-Risk / Service Issue",
        "  P ≥ 0.65 & low data usage  →  High-Risk / Low Usage",
        f"  P ≥ 0.65 & charges > K{HIGH_VALUE_KWACHA}  →  High-Risk / High Value",
        "  P ≥ 0.40                   →  At-Risk",
        "  tenure < 6 months          →  New Subscriber",
        "  high usage / charges       →  Loyal / High Value  |  else Stable / Standard",
    ])

    if best:
        _ppt_bullet_slide(prs, "Model Selection Rationale", [
            f"Selected model: {best}",
            f"Test ROC-AUC: {ss.model_results[best]['roc_auc']:.3f}",
            "Criteria: highest test ROC-AUC among six candidates.",
            "Production: batch score daily; expose probability + segment to CRM API.",
        ])

    _ppt_bullet_slide(prs, "Deployment Considerations", [
        "Retrain quarterly or when AUC drops > 5 points.",
        "Store model artefact + scaler + feature list versioned in MLflow.",
        "Log predictions with SHAP/feature contributions for explainability.",
        "Alert when segment distribution shifts beyond historical bounds.",
    ])

    if ss.eda_insights:
        _ppt_bullet_slide(prs, "EDA Key Findings", ss.eda_insights[:8])

    if ss.recommendations_df is not None:
        seg_detail = ss.recommendations_df.groupby("segment").agg(
            count=("churn_probability", "count"),
            avg_prob=("churn_probability", "mean"),
            max_prob=("churn_probability", "max"),
        ).reset_index()
        headers = ["Segment", "Count", "Avg P(churn)", "Max P(churn)"]
        rows = [
            [r["segment"], int(r["count"]), f"{r['avg_prob']:.3f}", f"{r['max_prob']:.3f}"]
            for _, r in seg_detail.iterrows()
        ]
        _ppt_table_slide(prs, "Segment Detail Table", headers, rows)

    _ppt_bullet_slide(prs, "Reproducibility & Governance", [
        "Random seed: 42 for split and tree-based models.",
        "Train/test split: 80/20 with stratification when class counts allow.",
        "Excluded identifiers: subscriber_id, msisdn, id_number, account_number.",
        "Excluded dates from encoding: activation_date, created_date, closed_date.",
        "Pipeline version exported as JSON alongside PowerPoint decks.",
    ])

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()






# ─────────────────────────────────────────────────────────────────────────────
# STEP 9 — DATA STORY
# ─────────────────────────────────────────────────────────────────────────────
def step_story() -> None:
    st.markdown(
        f"""
<div style="background:linear-gradient(135deg,#1565c0,#0d47a1);
    color:#fff;border-radius:16px;padding:28px 32px;margin-bottom:24px;">
  <h2 style="margin:0 0 6px 0;">📡 Zambia Telco Subscriber Analytics</h2>
  <p style="margin:0;opacity:.85;font-size:1rem;">
    Executive Data Story — {ss.dataset_name} &nbsp;·&nbsp; {datetime.now().strftime("%B %Y")}
    &nbsp;·&nbsp; Download the PowerPoint deck below for stakeholder presentations
  </p>
</div>
""",
        unsafe_allow_html=True,
    )

    if ss.raw_df is None:
        st.warning("Complete earlier steps first.")
        return

    df = ss.raw_df
    tgt = ss.target_col
    churn_rate = pd.to_numeric(df[tgt], errors="coerce").mean() if tgt in df.columns else None

    st.subheader("Executive KPIs")
    k1, k2, k3, k4, k5 = st.columns(5)
    with k1:
        _kpi("Subscribers", f"{df.shape[0]:,}")
    with k2:
        _kpi("Features", len(ss.feature_cols) if ss.feature_cols else df.shape[1], "green")
    with k3:
        _kpi("Churn Rate", f"{churn_rate:.1%}" if churn_rate is not None else "N/A", "orange")
    with k4:
        best = (
            max(ss.model_results, key=lambda n: ss.model_results[n]["roc_auc"])
            if ss.model_results else None
        )
        _kpi("Best ROC-AUC", f"{ss.model_results[best]['roc_auc']:.3f}" if best else "N/A", "purple")
    with k5:
        rar = ss.analytics_kpis.get("total_revenue_at_risk")
        _kpi("Revenue at Risk", _fmt_kwacha(rar) if rar else "N/A", "orange")

    st.subheader("Pipeline Journey")
    journey = []
    step_notes = {
        1: "CSV uploaded or demo data generated",
        2: "EDA insights captured" if ss.eda_insights else "Pending analysis",
        3: f"{len(ss.clean_log)} cleaning steps applied" if ss.clean_log else "Pending cleaning",
        4: f"{len(ss.feature_cols)} features selected" if ss.feature_cols else "Pending engineering",
        5: f"{len(ss.models)} models trained" if ss.models else "Pending training",
        6: "Evaluation metrics computed" if ss.model_results else "Pending evaluation",
        7: "Bias diagnostics available" if ss.models else "Pending bias check",
        8: "Segments assigned" if ss.recommendations_df is not None else "Pending recommendations",
        9: "Executive story ready" if 9 in ss.steps_done else "Pending export",
        10: "Market KPIs monitored" if 10 in ss.steps_done else "Pending market intelligence",
    }
    for num, icon, label in STEPS:
        status = "✅ Complete" if num in ss.steps_done else "○ Pending"
        journey.append({
            "Step": f"{num} {icon} {label}",
            "Status": status,
            "Notes": step_notes.get(num, ""),
        })
    st.dataframe(pd.DataFrame(journey), use_container_width=True, hide_index=True)

    st.subheader("Data Lineage")
    lineage_rows = [
        ("Raw ingestion", ss.raw_df.shape if ss.raw_df is not None else "—"),
        ("After cleaning", ss.cleaned_df.shape if ss.cleaned_df is not None else "—"),
        ("Feature matrix (train)", ss.X_train.shape if ss.X_train is not None else "—"),
        ("Feature matrix (test)", ss.X_test.shape if ss.X_test is not None else "—"),
        ("Recommendations scored", ss.recommendations_df.shape if ss.recommendations_df is not None else "—"),
    ]
    st.dataframe(
        pd.DataFrame(lineage_rows, columns=["Stage", "Shape (rows × cols)"]),
        use_container_width=True,
        hide_index=True,
    )

    if churn_rate is not None:
        st.subheader("Churn Status Assessment")
        if churn_rate > 0.25:
            _warn("Status: **CRITICAL** — churn above 25% threshold. Immediate retention programme required.")
        elif churn_rate > 0.15:
            _warn("Status: **ELEVATED** — proactive retention recommended.")
        else:
            _ok("Status: Within acceptable range.")

    if ss.clean_log:
        st.subheader("Cleaning Highlights")
        for entry in ss.clean_log[:6]:
            _ok(entry)

    if ss.eda_insights:
        st.subheader("Key EDA Insights")
        for insight in ss.eda_insights:
            _insight(insight)

    if _operator_col(df):
        st.subheader("Zambia Market KPI Summary")
        op_kpis = _compute_operator_kpis(df, tgt)
        _render_market_kpi_cards(op_kpis, ss.focus_operator or FOCUS_OPERATOR)
        st.dataframe(op_kpis.round(2), use_container_width=True, hide_index=True)

    if ss.model_results:
        st.subheader("Model Scorecard")
        summary = pd.DataFrame({
            name: {k: v for k, v in res.items() if k not in ("y_pred", "y_prob")}
            for name, res in ss.model_results.items()
        }).T.round(4)
        st.dataframe(summary, use_container_width=True)

        if ss.feature_importance is not None:
            st.subheader("Top 8 Churn Drivers (Random Forest)")
            fi = ss.feature_importance.head(8).reset_index()
            fi.columns = ["Feature", "Importance"]
            fig = px.bar(fi, x="Importance", y="Feature", orientation="h")
            st.plotly_chart(fig, use_container_width=True)

    if ss.recommendations_df is not None:
        st.subheader("Segment Snapshot")
        reco_story = ss.recommendations_df
        seg = reco_story.groupby("segment").agg(
            subscribers=("churn_probability", "count"),
            avg_prob=("churn_probability", "mean"),
            **({"total_rar": ("revenue_at_risk", "sum")} if "revenue_at_risk" in reco_story.columns else {}),
        ).reset_index()
        st.dataframe(seg.round(3), use_container_width=True, hide_index=True)
        hi = int((reco_story["churn_probability"] >= 0.5).sum())
        rar = ss.analytics_kpis.get("total_revenue_at_risk", 0)
        _insight(
            f"**{hi:,} subscribers** have P(churn) ≥ 50%. "
            f"Total revenue at risk: **{_fmt_kwacha(rar)}**."
        )

    if ss.secondary_results:
        st.subheader("Secondary Predictions")
        st.dataframe(pd.DataFrame(ss.secondary_results).T.round(4), use_container_width=True)

    st.subheader("Documentation & Methodology")
    with st.expander("Pipeline Methodology Reference", expanded=False):
        st.markdown(
            """
**Step 1 — Data Ingestion**  
Upload CSV/Excel, load 2,000-row demo data, or download the 500-row sample CSV.  
44 telco fields spanning identity, market (operator, report_month), subscription, usage,
network KPIs, CX, billing, and loyalty — localised for Zambia (ZMW, NRC, MSISDN).

**Step 2 — Exploratory Data Analysis**  
Descriptive statistics, target distribution (bar chart uses column name for colour),  
Pearson correlation heatmap, missing-value audit, numeric histograms, categorical churn rates.

**Step 3 — Data Cleaning**  
Duplicate removal, null-token normalisation (N/A, null, None), whitespace stripping,  
median imputation, numeric type coercion, re-imputation after coercion.

**Step 4 — Feature Engineering**  
Derived KPIs: `charge_per_mb`, `complaint_rate`, `engagement_score` (z-score mean only),  
`tenure_bucket`, `resolution_hours`, `sla_status`.  
Categorical encoding with `_SKIP_ENCODE_COLS` and `_MAX_OHE_CARDINALITY=30`.  
Binary target validation via `_prepare_binary_target`. Safe stratified split via `_safe_train_test_split`.  
Imputation and StandardScaler fit on train only.

**Step 5 — Model Training**  
Six classifiers: LogisticRegression, RandomForest, GradientBoosting, SVC(probability=True),  
KNeighborsClassifier, AdaBoostClassifier. NaN safety impute before train.

**Step 6 — Model Evaluation**  
Confusion matrix, ROC curves for all models, classification report, probability histogram.

**Step 7 — Bias Detection**  
Learning curves and train vs test AUC gap per model (overfit / underfit / good fit).

**Step 8 — Recommendations**  
Churn probability scoring, 7 segments, product + campaign mapping, CSV export.

**Step 9 — Data Story**  
Executive UI, `_generate_pptx` (executive deck), `_generate_technical_pptx` (technical deck), JSON export.
"""
        )

    st.subheader("Export Centre")
    st.markdown(
        "**Business audience:** KPIs, churn profile, scorecard, campaigns, action plan.  \n"
        "**Technical audience:** pipeline architecture, feature engineering, model zoo, bias diagnostics."
    )

    col_a, col_b, col_c = st.columns(3)
    with col_a:
        try:
            ppt_bytes = _generate_pptx()
            st.download_button(
                "📊 Download Executive Deck (.pptx)",
                data=ppt_bytes,
                file_name=f"telco_executive_{datetime.now().strftime('%Y%m%d')}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )
        except Exception as exc:
            st.error(f"Executive PPT generation failed: {exc}")

    with col_b:
        try:
            tech_bytes = _generate_technical_pptx()
            st.download_button(
                "🔬 Download Technical Deck (.pptx)",
                data=tech_bytes,
                file_name=f"telco_technical_{datetime.now().strftime('%Y%m%d')}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )
        except Exception as exc:
            st.error(f"Technical PPT generation failed: {exc}")

    with col_c:
        export_payload = {
            "generated_at": datetime.now().isoformat(),
            "dataset": ss.dataset_name,
            "target_column": ss.target_col,
            "rows": int(df.shape[0]) if df is not None else 0,
            "features_used": ss.feature_cols,
            "steps_completed": sorted(ss.steps_done),
            "eda_insights": ss.eda_insights,
            "clean_log": ss.clean_log,
            "model_results": {
                name: {k: float(v) if isinstance(v, (np.floating, float)) else v
                       for k, v in res.items() if k not in ("y_pred", "y_prob")}
                for name, res in ss.model_results.items()
            } if ss.model_results else {},
            "best_model": best,
            "feature_importance_top10": (
                ss.feature_importance.head(10).to_dict() if ss.feature_importance is not None else {}
            ),
            "segment_summary": (
                ss.recommendations_df.groupby("segment")["churn_probability"].agg(["count", "mean"]).to_dict()
                if ss.recommendations_df is not None else {}
            ),
        }
        st.download_button(
            "📄 Download JSON Report",
            data=json.dumps(export_payload, indent=2, default=str).encode(),
            file_name=f"telco_pipeline_report_{datetime.now().strftime('%Y%m%d')}.json",
            mime="application/json",
            use_container_width=True,
        )

    if st.button("Proceed to Market Intelligence →", type="primary"):
        _complete_step(9)
        st.rerun()


# ─────────────────────────────────────────────────────────────────────────────
# STEP 10 — MARKET INTELLIGENCE
# ─────────────────────────────────────────────────────────────────────────────
def step_market() -> None:
    _step_header(
        10, "Market Intelligence",
        "Monitor key KPIs (subscriber base, revenue, ARPU, churn, NPS) and benchmark "
        f"**{FOCUS_OPERATOR}** against **MTN Zambia** and **Airtel Zambia** — "
        "portfolio metrics plus national market reference gaps.",
    )
    if ss.raw_df is None:
        st.warning("Complete Step 1 first.")
        return

    df = ss.raw_df.copy()
    tgt = ss.target_col
    op_kpis = _compute_operator_kpis(df, tgt)
    ss.market_kpis = op_kpis.to_dict("records") if not op_kpis.empty else {}

    st.subheader("KPI Monitor")
    focus_opts = list(op_kpis["operator"]) if not op_kpis.empty else list(ZAMBIA_OPERATORS)
    focus = st.selectbox(
        "Focus operator for KPI cards",
        focus_opts,
        index=focus_opts.index(ss.focus_operator) if ss.focus_operator in focus_opts else 0,
        key="market_step_focus",
    )
    ss.focus_operator = focus
    _render_market_kpi_cards(op_kpis, focus)

    total_subs = int(op_kpis["subscribers"].sum()) if not op_kpis.empty else len(df)
    total_rev = float(op_kpis["monthly_revenue"].sum()) if not op_kpis.empty else 0
    m1, m2, m3, m4 = st.columns(4)
    m1.metric("Total Portfolio Subscribers", f"{total_subs:,}")
    m2.metric("Total Monthly Revenue", _fmt_kwacha(total_rev))
    if not op_kpis.empty and tgt in df.columns:
        m3.metric("Portfolio Churn Rate", f"{op_kpis['churn_rate'].mean() * 100:.1f}%")
    if "report_month" in df.columns:
        m4.metric("Reporting Periods", df["report_month"].nunique())

    _render_operator_benchmark_charts(df, tgt)

    st.subheader("KPI Alert Thresholds")
    if not op_kpis.empty and FOCUS_OPERATOR in op_kpis["operator"].values:
        zam = op_kpis.loc[op_kpis["operator"] == FOCUS_OPERATOR].iloc[0]
        ref = ZAMBIA_MARKET_REFERENCE.set_index("operator").loc[FOCUS_OPERATOR]
        alerts = []
        if zam.get("churn_rate_pct", 0) > ref["churn_rate_pct"]:
            alerts.append(f"Churn above national reference ({zam['churn_rate_pct']:.1f}% vs {ref['churn_rate_pct']:.1f}%)")
        if zam.get("avg_arpu", 0) < ref["monthly_arpu_kwacha"]:
            alerts.append(f"ARPU below national reference ({_fmt_kwacha(zam['avg_arpu'])} vs {_fmt_kwacha(ref['monthly_arpu_kwacha'])})")
        if zam.get("avg_nps", 100) < ref["nps"]:
            alerts.append(f"NPS below national reference ({zam['avg_nps']:.0f} vs {ref['nps']:.0f})")
        if alerts:
            for a in alerts:
                _warn(a)
        else:
            _ok(f"{FOCUS_OPERATOR} is within national reference thresholds on monitored KPIs.")

    if st.button("Mark Pipeline Complete ✅", type="primary"):
        _complete_step(10)
        st.balloons()
        st.success("🎉 Zambia Telco Analytics Pipeline complete!")


# ─────────────────────────────────────────────────────────────────────────────
# ROUTER
# ─────────────────────────────────────────────────────────────────────────────
STEP_FNS = {
    1: step_ingest,
    2: step_eda,
    3: step_clean,
    4: step_features,
    5: step_train,
    6: step_evaluate,
    7: step_bias,
    8: step_recommendations,
    9: step_story,
    10: step_market,
}

if __name__ == "__main__":
    _nav_sidebar()
    STEP_FNS[ss.current_step if ss.current_step in STEP_FNS else len(STEPS)]()
